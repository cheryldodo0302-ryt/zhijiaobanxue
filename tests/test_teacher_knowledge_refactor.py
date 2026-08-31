import io
import json
from pathlib import Path

from docx import Document
from pptx import Presentation

from auth_service import AuthService
from campus_service import CampusService
from database import LearningDatabase
from formula_client import Pix2TextClient
from ingestion_service import IngestionService
from job_secret_store import decrypt_job_secret
from mineru_client import MinerUClient
from teacher_service import TeacherService


def teacher_scope(tmp_path: Path):
    db = LearningDatabase(tmp_path / "knowledge-refactor.db")
    campus = CampusService(db, tmp_path / "uploads", provider_factory=lambda: None)
    teacher = AuthService(db, tmp_path / "secret").create_user(
        "teacher", "safe-password-123", "teacher"
    )
    course = TeacherService(db, campus).create_course(teacher, "数据库")
    return db, campus, teacher, course


class ForbiddenMinerU:
    enabled = True

    def parse(self, *_args, **_kwargs):
        raise AssertionError("native text/Office content must not invoke MinerU")


class DisabledFormula:
    enabled = False


def presentation_with_repeated_and_ambiguous_titles() -> bytes:
    presentation = Presentation()
    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = "数据库课程"
    cover.placeholders[1].text = "课程封面"
    for title, body in (
        ("1.2 关系模型", "实体与关系"),
        ("1.2 关系模型", "关系的完整性"),
        ("1.3 关系代数", "选择与投影"),
        ("3 第三点", "这是 1.3 下的第三条展开，不是第三章"),
    ):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue()


class PptHierarchyProvider:
    def generate(self, _system, prompt):
        payload = json.loads(prompt)
        if "presentation_groups" in payload:
            groups = payload["presentation_groups"]
            assert [group["title"] for group in groups] == [
                "1.2 关系模型", "1.3 关系代数", "3 第三点",
            ]
            return json.dumps({"outline": [{
                "group_id": group["group_id"],
                "hierarchy_level": "point",
                "chapter": "第 1 章",
                "section": "1.3 关系代数" if group["title"].startswith("3") else group["title"],
                "reason": "结合相邻标题与正文判定",
            } for group in groups]}, ensure_ascii=False)
        if "source_points" in payload:
            # Deliberately reverse the model response. The service must restore
            # source/PPT order before materializing the course tree.
            points = list(reversed(payload["source_points"]))
            return json.dumps({"points": [{
                "course_key": f"course-{index}", "chapter": item["chapter"],
                "section": item["section"], "title": item["title"],
                "keywords": [], "source_node_ids": [item["source_node_id"]],
            } for index, item in enumerate(points)], "relations": []}, ensure_ascii=False)
        raise AssertionError(f"unexpected PPT semantic payload: {payload.keys()}")


class RecordingMinerU:
    enabled = True

    def __init__(self):
        self.method = ""

    def parse(self, _path, *, method, asset_dir):
        self.method = method
        return {
            "_version_name": "test",
            "_markdown": "# 第一章\n\n关系模型定义",
            "_image_paths": {},
            "pdf_info": [{
                "page_idx": 0,
                "para_blocks": [
                    {"type": "title", "content": "# 第一章"},
                    {"type": "text", "content": "关系模型定义"},
                ],
            }],
        }


def test_txt_uses_native_markdown_and_inline_text_preview(tmp_path):
    db, campus, teacher, course = teacher_scope(tmp_path)
    service = IngestionService(db, campus)
    service.mineru = ForbiddenMinerU()
    body = "第一节\n关系模型定义"
    job = service.queue_document(
        teacher, course["course_id"], "lesson.txt", "text/plain", body.encode()
    )
    service.process_job(job["job_id"])
    artifact = db.fetch_one(
        "SELECT * FROM document_artifacts WHERE document_id=? AND artifact_type='canonical_markdown'",
        (job["document_id"],),
    )
    assert artifact and artifact["status"] == "ready"
    assert Path(artifact["stored_path"]).read_text(encoding="utf-8").strip() == body
    assert service.preview_descriptor(teacher, job["document_id"])["preview_kind"] == "text"
    media_type, preview = service.preview_file(teacher, job["document_id"])
    assert media_type == "text/plain" and preview == body


def test_pptx_uses_native_text_and_supports_browser_preview_without_converter(tmp_path, monkeypatch):
    db, campus, teacher, course = teacher_scope(tmp_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "关系模型"
    slide.placeholders[1].text = "实体与关系"
    stream = io.BytesIO()
    presentation.save(stream)
    service = IngestionService(db, campus)
    service.mineru = ForbiddenMinerU()
    monkeypatch.setattr("ingestion_service.shutil.which", lambda _name: None)
    job = service.queue_document(
        teacher, course["course_id"], "lesson.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation", stream.getvalue(),
    )
    service.process_job(job["job_id"])
    descriptor = service.preview_descriptor(teacher, job["document_id"])
    assert descriptor["preview_kind"] == "pptx"
    assert descriptor["conversion_status"] == "ready"
    assert descriptor["total_pages"] == 1
    assert db.fetch_one(
        "SELECT 1 ok FROM document_artifacts WHERE document_id=? AND artifact_type='canonical_markdown' AND status='ready'",
        (job["document_id"],),
    )


def test_docx_supports_safe_browser_preview_without_libreoffice(tmp_path, monkeypatch):
    _, campus, teacher, course = teacher_scope(tmp_path)
    document = Document()
    document.add_heading("实验九 信息系统分析与设计", level=1)
    document.add_paragraph("目标与要求")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "教学内容"
    table.cell(0, 1).text = "完成业务流程建模"
    stream = io.BytesIO()
    document.save(stream)
    service = IngestionService(campus.db, campus)
    service.mineru = ForbiddenMinerU()
    monkeypatch.setattr("ingestion_service.shutil.which", lambda _name: None)
    job = service.queue_document(
        teacher, course["course_id"], "outline.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        stream.getvalue(),
    )
    service.process_job(job["job_id"])

    descriptor = service.preview_descriptor(teacher, job["document_id"])
    assert descriptor["preview_kind"] == "docx"
    assert descriptor["conversion_status"] == "ready"
    media_type, preview = service.preview_file(teacher, job["document_id"])
    assert media_type == "text/html"
    assert "实验九 信息系统分析与设计" in preview
    assert "完成业务流程建模" in preview
    assert "<table>" in preview


def test_pptx_groups_consecutive_same_titles_and_exposes_slide_numbers(tmp_path, monkeypatch):
    db, campus, teacher, course = teacher_scope(tmp_path)
    service = IngestionService(db, campus)
    service.mineru = ForbiddenMinerU()
    monkeypatch.setattr("ingestion_service.shutil.which", lambda _name: None)
    job = service.queue_document(
        teacher, course["course_id"], "ordered.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        presentation_with_repeated_and_ambiguous_titles(),
    )
    service.process_job(job["job_id"])

    candidates = service.list_knowledge_candidates(teacher, job["document_id"])
    assert [candidate["title"] for candidate in candidates] == [
        "1.2 关系模型", "1.3 关系代数", "3 第三点",
    ]
    assert {block["page_number"] for block in candidates[0]["source_blocks"]} == {2, 3}
    assert "实体与关系" in candidates[0]["source_markdown"]
    assert "关系的完整性" in candidates[0]["source_markdown"]
    slides = service.list_presentation_slides(teacher, job["document_id"])
    assert [slide["slide_number"] for slide in slides] == [1, 2, 3, 4, 5]
    assert slides[1]["title"] == "1.2 关系模型"
    assert slides[1]["regions"][0]["role"] == "title"


def test_ppt_title_rebuild_is_local_and_does_not_enqueue_semantic_api(tmp_path, monkeypatch):
    db, campus, teacher, course = teacher_scope(tmp_path)
    service = IngestionService(db, campus)
    service.mineru = ForbiddenMinerU()
    monkeypatch.setattr("ingestion_service.shutil.which", lambda _name: None)
    job = service.queue_document(
        teacher, course["course_id"], "第一章 数据库系统概述.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        presentation_with_repeated_and_ambiguous_titles(),
    )
    service.process_job(job["job_id"])
    assert db.fetch_one(
        "SELECT COUNT(*) n FROM semantic_analysis_jobs WHERE document_id=?", (job["document_id"],)
    )["n"] == 1

    rebuilt = service.reparse_presentation(teacher, job["job_id"])
    assert rebuilt["pipeline_stage"] == "PPT_TITLE_REBUILD"
    service.process_job(job["job_id"])

    assert db.fetch_one(
        "SELECT COUNT(*) n FROM semantic_analysis_jobs WHERE document_id=?", (job["document_id"],)
    )["n"] == 1
    assert [row["title"] for row in service.list_knowledge_candidates(teacher, job["document_id"])] == [
        "1.2 关系模型", "1.3 关系代数", "3 第三点",
    ]
    document_outline = service.document_outline(teacher, job["document_id"])
    document_points = [
        row for row in document_outline["nodes"] if row["node_type"] == "knowledge_point"
    ]
    assert [row["title"] for row in document_points] == [
        "1.2 关系模型", "1.3 关系代数", "3 第三点",
    ]
    assert document_points[0]["source_pages"] == [2, 3]
    assert "实体与关系" in document_points[0]["markdown"]
    assert "关系的完整性" in document_points[0]["markdown"]
    course_outline = service.course_outline(teacher, course["course_id"], material_type="slides")
    assert [
        row["title"] for row in course_outline["nodes"] if row["node_type"] == "chapter"
    ] == ["第一章 数据库系统概述"]
    assert [
        row["title"] for row in course_outline["nodes"]
        if row["node_type"] == "knowledge_point"
    ] == ["1.2 关系模型", "1.3 关系代数", "3 第三点"]


def test_candidate_and_document_leaf_reviews_are_bidirectionally_synchronized(
    tmp_path, monkeypatch
):
    db, campus, teacher, course = teacher_scope(tmp_path)
    service = IngestionService(db, campus)
    service.mineru = ForbiddenMinerU()
    monkeypatch.setattr("ingestion_service.shutil.which", lambda _name: None)
    job = service.queue_document(
        teacher, course["course_id"], "sync.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        presentation_with_repeated_and_ambiguous_titles(),
    )
    service.process_job(job["job_id"])
    service.reparse_presentation(teacher, job["job_id"])
    service.process_job(job["job_id"])

    candidates = service.list_knowledge_candidates(teacher, job["document_id"])
    assert len(candidates) == 3
    assert all(candidate.get("document_node_id") for candidate in candidates)

    approved = service.approve_knowledge_candidate(
        teacher, candidates[0]["candidate_id"]
    )
    first_node = db.fetch_one(
        "SELECT status FROM knowledge_nodes WHERE node_id=?",
        (candidates[0]["document_node_id"],),
    )
    assert approved["review_status"] == "APPROVED"
    assert first_node["status"] == "approved"

    service.update_node(
        teacher, candidates[1]["document_node_id"], {"status": "approved"}
    )
    refreshed = service.list_knowledge_candidates(teacher, job["document_id"])
    assert refreshed[1]["review_status"] == "APPROVED"

    third_node = db.fetch_one(
        "SELECT parent_id FROM knowledge_nodes WHERE node_id=?",
        (candidates[2]["document_node_id"],),
    )
    result = service.approve_nodes_batch(teacher, [third_node["parent_id"]])
    assert result["approved_leaf_count"] == 1
    refreshed = service.list_knowledge_candidates(teacher, job["document_id"])
    assert {candidate["review_status"] for candidate in refreshed} == {"APPROVED"}


def test_parent_approval_recursively_approves_all_children_and_candidate_queue_has_body(
    tmp_path, monkeypatch
):
    db, campus, teacher, course = teacher_scope(tmp_path)
    service = IngestionService(db, campus)
    service.mineru = ForbiddenMinerU()
    monkeypatch.setattr("ingestion_service.shutil.which", lambda _name: None)
    job = service.queue_document(
        teacher, course["course_id"], "complete-sections.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        presentation_with_repeated_and_ambiguous_titles(),
    )
    service.process_job(job["job_id"])
    service.reparse_presentation(teacher, job["job_id"])
    service.process_job(job["job_id"])

    candidates = service.list_knowledge_candidates(teacher, job["document_id"])
    assert candidates
    assert all(candidate["section_markdown"].strip() for candidate in candidates)
    assert all(candidate["section_excerpt"].strip() for candidate in candidates)

    first_node_id = candidates[0]["document_node_id"]
    section_id = db.fetch_one(
        "SELECT parent_id FROM knowledge_nodes WHERE node_id=?", (first_node_id,)
    )["parent_id"]
    chapter_id = db.fetch_one(
        "SELECT parent_id FROM knowledge_nodes WHERE node_id=?", (section_id,)
    )["parent_id"]
    service.update_node(teacher, chapter_id, {"status": "approved"})
    descendants = db.fetch_all(
        """WITH RECURSIVE tree(node_id) AS (
               SELECT node_id FROM knowledge_nodes WHERE node_id=?
               UNION ALL SELECT n.node_id FROM knowledge_nodes n JOIN tree t ON n.parent_id=t.node_id
           ) SELECT n.status FROM knowledge_nodes n JOIN tree USING(node_id)""",
        (chapter_id,),
    )
    assert descendants
    assert {row["status"] for row in descendants} == {"approved"}
    refreshed = service.list_knowledge_candidates(teacher, job["document_id"])
    approved_in_branch = [
        candidate for candidate in refreshed
        if chapter_id in {
            db.fetch_one(
                "SELECT parent_id FROM knowledge_nodes WHERE node_id=(SELECT parent_id FROM knowledge_nodes WHERE node_id=?)",
                (candidate["document_node_id"],),
            )["parent_id"]
        }
    ]
    assert approved_in_branch
    assert {candidate["review_status"] for candidate in approved_in_branch} == {"APPROVED"}


def test_heading_only_document_leaf_is_not_a_candidate_queue_item(tmp_path, monkeypatch):
    db, campus, teacher, course = teacher_scope(tmp_path)
    service = IngestionService(db, campus)
    service.mineru = ForbiddenMinerU()
    monkeypatch.setattr("ingestion_service.shutil.which", lambda _name: None)
    job = service.queue_document(
        teacher, course["course_id"], "heading-placeholder.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        presentation_with_repeated_and_ambiguous_titles(),
    )
    service.process_job(job["job_id"])
    service.reparse_presentation(teacher, job["job_id"])
    service.process_job(job["job_id"])
    before = service.list_knowledge_candidates(teacher, job["document_id"])
    placeholder = before[0]
    db.execute(
        "UPDATE knowledge_nodes SET markdown=? WHERE node_id=?",
        (f"# {placeholder['title']}", placeholder["document_node_id"]),
    )
    after = service.list_knowledge_candidates(teacher, job["document_id"])
    assert placeholder["candidate_id"] not in {row["candidate_id"] for row in after}


def test_syllabus_teaching_information_is_archived_not_exposed_as_knowledge(tmp_path):
    db, campus, teacher, course = teacher_scope(tmp_path)
    service = IngestionService(db, campus)
    document = campus.upload_document(
        course["course_id"], teacher["user_id"], teacher["role"],
        "课程教学大纲.txt", "text/plain", b"placeholder",
    )
    db.execute(
        """INSERT INTO document_material_metadata(
               document_id,material_type,suggested_material_type,classification_status
           ) VALUES(?,'syllabus','syllabus','confirmed')""",
        (document["document_id"],),
    )
    chapter_id, section_id = "kn_archive_chapter", "kn_archive_section"
    db.execute(
        """INSERT INTO knowledge_nodes(
               node_id,course_id,document_id,node_scope,node_type,title,sort_order,material_type
           ) VALUES(?,?,?,'document','chapter','一、课程基本信息',1,'syllabus')""",
        (chapter_id, course["course_id"], document["document_id"]),
    )
    db.execute(
        """INSERT INTO knowledge_nodes(
               node_id,course_id,document_id,node_scope,parent_id,node_type,title,sort_order,material_type
           ) VALUES(?,?,?,'document',?,'section','培养目标',2,'syllabus')""",
        (section_id, course["course_id"], document["document_id"], chapter_id),
    )
    db.execute(
        """INSERT INTO knowledge_nodes(
               node_id,course_id,document_id,node_scope,parent_id,node_type,title,markdown,
               sort_order,material_type
           ) VALUES(?,?,?,'document',?,'knowledge_point','培养目标','培养学生掌握数据库设计能力',3,'syllabus')""",
        ("kn_objective", course["course_id"], document["document_id"], section_id),
    )
    knowledge_chapter, knowledge_section = "kn_lab_chapter", "kn_lab_section"
    db.execute(
        """INSERT INTO knowledge_nodes(
               node_id,course_id,document_id,node_scope,node_type,title,sort_order,material_type
           ) VALUES(?,?,?,'document','chapter','四、实验教学内容纲要',4,'syllabus')""",
        (knowledge_chapter, course["course_id"], document["document_id"]),
    )
    db.execute(
        """INSERT INTO knowledge_nodes(
               node_id,course_id,document_id,node_scope,parent_id,node_type,title,sort_order,material_type
           ) VALUES(?,?,?,'document',?,'section','实验四',5,'syllabus')""",
        (knowledge_section, course["course_id"], document["document_id"], knowledge_chapter),
    )
    db.execute(
        """INSERT INTO knowledge_nodes(
               node_id,course_id,document_id,node_scope,parent_id,node_type,title,markdown,
               sort_order,material_type
           ) VALUES(?,?,?,'document',?,'knowledge_point','SQL程序设计（三）','掌握复杂SQL查询',6,'syllabus')""",
        ("kn_lab", course["course_id"], document["document_id"], knowledge_section),
    )

    archive = service.teaching_archive(teacher, course["course_id"])
    assert [row["title"] for row in archive["sections"]] == ["培养目标"]
    assert archive["sections"][0]["teaching_category"] == "objectives"
    outline = service.document_outline(teacher, document["document_id"])
    assert "培养目标" not in {row["title"] for row in outline["nodes"]}
    assert "SQL程序设计（三）" in {row["title"] for row in outline["nodes"]}


def test_teaching_archive_upload_can_target_one_or_multiple_classes(tmp_path):
    db, campus, teacher, course = teacher_scope(tmp_path)
    teachers = TeacherService(db, campus)
    term = teachers.create_term(
        teacher, "2026-2027 秋季", academic_year="2026-2027", teaching_period="秋季学期"
    )
    class_a = teachers.create_class(
        teacher, course["course_id"], term["term_id"], "数据库A班", "A班（本部）"
    )
    class_b = teachers.create_class(
        teacher, course["course_id"], term["term_id"], "数据库B班", "B班（仁济）"
    )
    service = IngestionService(db, campus)
    job = service.queue_teaching_archive_document_stream(
        teacher, course["course_id"], "A-B班教学大纲.txt", "text/plain",
        io.BytesIO("课程目标与考核方式".encode()),
        [class_a["class_id"], class_b["class_id"]], analysis_mode="local",
    )
    assignments = db.fetch_all(
        """SELECT class_id FROM teaching_archive_document_assignments
           WHERE document_id=? ORDER BY class_id""",
        (job["document_id"],),
    )
    assert {row["class_id"] for row in assignments} == {
        class_a["class_id"], class_b["class_id"],
    }
    metadata = db.fetch_one(
        "SELECT material_type,classification_status FROM document_material_metadata WHERE document_id=?",
        (job["document_id"],),
    )
    assert metadata == {"material_type": "syllabus", "classification_status": "confirmed"}
    archive = service.teaching_archive(teacher, course["course_id"])
    uploaded = next(row for row in archive["documents"] if row["document_id"] == job["document_id"])
    assert set(uploaded["class_labels"]) == {"A班（本部）", "B班（仁济）"}

    chapter_id, section_id, point_id = "kn_level_chapter", "kn_level_section", "kn_level_point"
    db.execute(
        """INSERT INTO knowledge_nodes(
               node_id,course_id,document_id,node_scope,node_type,title,sort_order,material_type
           ) VALUES(?,?,?,'document','chapter','实验教学内容',1,'syllabus')""",
        (chapter_id, course["course_id"], job["document_id"]),
    )
    db.execute(
        """INSERT INTO knowledge_nodes(
               node_id,course_id,document_id,node_scope,parent_id,node_type,title,sort_order,material_type
           ) VALUES(?,?,?,'document',?,'section','实验四',2,'syllabus')""",
        (section_id, course["course_id"], job["document_id"], chapter_id),
    )
    db.execute(
        """INSERT INTO knowledge_nodes(
               node_id,course_id,document_id,node_scope,parent_id,node_type,title,markdown,
               sort_order,material_type
           ) VALUES(?,?,?,'document',?,'knowledge_point','复杂SQL查询','完整知识正文',3,'syllabus')""",
        (point_id, course["course_id"], job["document_id"], section_id),
    )
    inherited = service.document_outline(teacher, job["document_id"])
    inherited_point = next(row for row in inherited["nodes"] if row["node_id"] == point_id)
    assert set(inherited_point["class_ids"]) == {class_a["class_id"], class_b["class_id"]}

    service.update_node(teacher, point_id, {"class_ids": [class_a["class_id"]]})
    class_a_outline = service.document_outline(
        teacher, job["document_id"], class_id=class_a["class_id"]
    )
    class_b_outline = service.document_outline(
        teacher, job["document_id"], class_id=class_b["class_id"]
    )
    assert point_id in {row["node_id"] for row in class_a_outline["nodes"]}
    assert point_id not in {row["node_id"] for row in class_b_outline["nodes"]}

    service.update_node(teacher, point_id, {"class_ids": []})
    course_wide = service.document_outline(
        teacher, job["document_id"], class_id="course_wide"
    )
    assert point_id in {row["node_id"] for row in course_wide["nodes"]}


def test_teacher_ai_settings_are_encrypted_and_reusable(tmp_path, monkeypatch):
    db, campus, teacher, course = teacher_scope(tmp_path)
    monkeypatch.setattr("job_secret_store.KEY_PATH", tmp_path / "saved-ai-secret.key")
    service = IngestionService(db, campus)
    saved = service.save_teacher_ai_settings(
        teacher, provider="openai_compatible", base_url="https://example.com/v1",
        model="teacher-model", api_key="teacher-saved-secret",
    )
    assert saved["has_api_key"] is True
    assert saved["verification_status"] == "untested"
    assert "api_key" not in saved
    stored_setting = db.fetch_one(
        "SELECT api_key_encrypted FROM teacher_ai_settings WHERE teacher_id=?", (teacher["user_id"],)
    )
    assert stored_setting["api_key_encrypted"] != "teacher-saved-secret"
    assert decrypt_job_secret(stored_setting["api_key_encrypted"]) == "teacher-saved-secret"

    job = service.queue_document(
        teacher, course["course_id"], "saved-key.md", "text/markdown",
        "# 标题\n\n正文".encode("utf-8"), analysis_mode="local",
    )
    service.process_job(job["job_id"])
    analysis = service.queue_semantic_analysis(
        teacher, job["document_id"], analysis_mode="api",
        ai_settings={"use_saved": True},
    )
    encrypted = db.fetch_one(
        "SELECT ai_key_encrypted FROM semantic_analysis_jobs WHERE analysis_job_id=?",
        (analysis["analysis_job_id"],),
    )["ai_key_encrypted"]
    assert decrypt_job_secret(encrypted) == "teacher-saved-secret"
    assert "ai_key_encrypted" not in analysis


def test_ppt_semantic_analysis_only_resolves_hierarchy_and_keeps_slide_order(tmp_path, monkeypatch):
    db = LearningDatabase(tmp_path / "ppt-semantic.db")
    campus = CampusService(db, tmp_path / "uploads", provider_factory=PptHierarchyProvider)
    teacher = AuthService(db, tmp_path / "secret").create_user(
        "teacher", "safe-password-123", "teacher"
    )
    course = TeacherService(db, campus).create_course(teacher, "数据库")
    service = IngestionService(db, campus)
    service.mineru = ForbiddenMinerU()
    monkeypatch.setattr("ingestion_service.shutil.which", lambda _name: None)
    job = service.queue_document(
        teacher, course["course_id"], "semantic-order.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        presentation_with_repeated_and_ambiguous_titles(),
    )
    service.process_job(job["job_id"])
    analysis = db.fetch_one(
        "SELECT * FROM semantic_analysis_jobs WHERE document_id=?", (job["document_id"],)
    )
    service.process_semantic_analysis(analysis["analysis_job_id"])

    completed = service.get_analysis_job(teacher, analysis["analysis_job_id"])
    assert completed["status"] == "review_required", completed["error_message"]
    assert completed["analyzer_version"] == "ppt-title-outline-v1"
    assert completed["api_calls"] == 2
    document = service.document_outline(teacher, job["document_id"])
    course_outline = service.course_outline(teacher, course["course_id"])
    assert [node["title"] for node in document["nodes"] if node["node_type"] == "knowledge_point"] == [
        "1.2 关系模型", "1.3 关系代数", "3 第三点",
    ]
    assert [node["title"] for node in course_outline["nodes"] if node["node_type"] == "knowledge_point"] == [
        "1.2 关系模型", "1.3 关系代数", "3 第三点",
    ]
    candidates = service.list_knowledge_candidates(teacher, job["document_id"])
    assert [candidate["title"] for candidate in candidates] == [
        "1.2 关系模型", "1.3 关系代数", "3 第三点",
    ]
    assert {candidate["review_status"] for candidate in candidates} == {"PENDING"}
    assert all(candidate["chapter_path"] for candidate in candidates)


def test_ppt_semantic_retry_backfills_legacy_document_ir_titles_locally(tmp_path, monkeypatch):
    db = LearningDatabase(tmp_path / "legacy-ppt-semantic.db")
    campus = CampusService(db, tmp_path / "uploads", provider_factory=PptHierarchyProvider)
    teacher = AuthService(db, tmp_path / "secret").create_user(
        "teacher", "safe-password-123", "teacher"
    )
    course = TeacherService(db, campus).create_course(teacher, "数据库")
    service = IngestionService(db, campus)
    service.mineru = ForbiddenMinerU()
    monkeypatch.setattr("ingestion_service.shutil.which", lambda _name: None)
    job = service.queue_document(
        teacher, course["course_id"], "legacy-semantic-order.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        presentation_with_repeated_and_ambiguous_titles(),
    )
    service.process_job(job["job_id"])
    # Simulate a DocumentIR produced before PPT title metadata was persisted.
    db.execute(
        "UPDATE document_blocks SET raw_payload_json='{}' WHERE document_id=?",
        (job["document_id"],),
    )
    analysis = db.fetch_one(
        "SELECT * FROM semantic_analysis_jobs WHERE document_id=?", (job["document_id"],)
    )
    service.process_semantic_analysis(analysis["analysis_job_id"])

    completed = service.get_analysis_job(teacher, analysis["analysis_job_id"])
    assert completed["status"] == "review_required", completed["error_message"]
    result = json.loads(completed["result_json"])
    assert result["ppt_metadata_source"] == "local_ppt_title_refresh"
    assert [
        node["title"] for node in service.document_outline(teacher, job["document_id"])["nodes"]
        if node["node_type"] == "knowledge_point"
    ] == ["1.2 关系模型", "1.3 关系代数", "3 第三点"]


def test_teacher_pdf_uses_mineru_auto_and_persists_returned_markdown(tmp_path):
    db, campus, teacher, course = teacher_scope(tmp_path)
    service = IngestionService(db, campus)
    mineru = RecordingMinerU()
    service.mineru = mineru
    service.formula = DisabledFormula()
    job = service.queue_document(
        teacher, course["course_id"], "lesson.pdf", "application/pdf", b"%PDF-1.4\npdf fixture"
    )
    service.process_job(job["job_id"])
    assert mineru.method == "auto"
    artifact = db.fetch_one(
        "SELECT * FROM document_artifacts WHERE document_id=? AND artifact_type='canonical_markdown'",
        (job["document_id"],),
    )
    assert artifact and "关系模型定义" in Path(artifact["stored_path"]).read_text(encoding="utf-8")
    blocks = db.fetch_all("SELECT * FROM document_blocks WHERE document_id=?", (job["document_id"],))
    assert {row["parser_name"] for row in blocks} == {"MinerU"}


def test_remote_clients_attach_bearer_tokens_and_verify_tls(monkeypatch):
    monkeypatch.setenv("ZHIJIAO_MINERU_TOKEN", "mineru-secret")
    monkeypatch.setenv("ZHIJIAO_FORMULA_TOKEN", "formula-secret")
    monkeypatch.setenv("ZHIJIAO_MINERU_VERIFY_TLS", "1")
    monkeypatch.setenv("ZHIJIAO_FORMULA_VERIFY_TLS", "1")
    mineru = MinerUClient("https://mineru.example.edu")
    formula = Pix2TextClient("https://formula.example.edu")
    assert mineru.headers == {"Authorization": "Bearer mineru-secret"}
    assert formula.headers == {"Authorization": "Bearer formula-secret"}
    assert mineru.verify_tls is True and formula.verify_tls is True
