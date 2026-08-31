import base64
import io
import json
import zipfile
from unittest.mock import patch

import pytest

from agent_service import CampusAgentService
from campus_service import CampusService, NotFound, PermissionDenied, ValidationError
from database import LearningDatabase
from llm_provider import LLMProvider, QwenProvider
from skills.memory.service import _normalize_judgment_answer, _normalize_question_item


class StubProvider(LLMProvider):
    def generate(self, system_prompt: str, user_prompt: str) -> str:
        if "语义分块专家" in system_prompt:
            return json.dumps([{"title":"《监督学习的2个要点》","keywords":["监督学习","标签"],
                                "content":"监督学习使用带标签样本训练模型，并根据误差更新参数。"}], ensure_ascii=False)
        if "背诵监督员" in system_prompt:
            return json.dumps({"score":80,"missing_points":["更新参数"],"error_points":[],"feedback":"补充参数更新过程。"}, ensure_ascii=False)
        if "题库结构化整理智能体" in system_prompt:
            return json.dumps([
                {"type":"单选题","question":"监督学习使用什么样本？","options":["带标签样本","无标签样本"],
                 "answer":"带标签样本","explanation":"题库答案","knowledge_point":"监督学习"},
                {"type":"简答题","question":"说明监督学习的含义。","options":[],"answer":"",
                 "explanation":"","knowledge_point":"监督学习"},
            ], ensure_ascii=False)
        if "出题专家" in system_prompt:
            return json.dumps([{"type":"简答题","question":"监督学习使用什么样本？","options":[],
                                "answer":"带标签样本","explanation":"材料原文。","knowledge_point":"监督学习"}], ensure_ascii=False)
        if "练习批改智能体" in system_prompt:
            if "错误答案" in user_prompt:
                return json.dumps({"score":0,"results":[{"index":1,"correct":False,"correct_answer":"带标签样本",
                                    "feedback":"应回答带标签样本"}],"weak_points":["监督学习"],
                                   "summary":"需要复习监督学习样本类型"}, ensure_ascii=False)
            return json.dumps({"score":100,"results":[{"index":1,"correct":True,"correct_answer":"带标签样本",
                                "feedback":"回答正确"}],"weak_points":[],"summary":"掌握良好"}, ensure_ascii=False)
        evidence = user_prompt.split("【课程证据】", 1)[-1]
        lines = [line.strip() for line in evidence.splitlines()
                 if line.strip() and not line.startswith("来源")]
        return "根据测试课程资料，" + " ".join(lines[:3])[:500]


@pytest.fixture()
def campus(tmp_path):
    return CampusService(LearningDatabase(tmp_path / "campus.db"), tmp_path / "uploads", StubProvider)


def add_shared(campus):
    course = campus.create_course("共享课", "shared_course", "teacher_1", "teacher", visibility="enrolled")
    campus.enroll_student(course["course_id"], "teacher_1", "student_1")
    campus.upload_document(course["course_id"], "teacher_1", "teacher", "lesson.md", "text/markdown",
                           "# 监督学习\n监督学习使用带标签样本训练模型。".encode())
    return course


def test_personal_course_and_shared_course_permissions(campus):
    personal = campus.create_course("我的课", "personal_course", "student_1", "student")
    with pytest.raises(PermissionDenied):
        campus.require_access(personal["course_id"], "teacher_1", "teacher")
    shared = add_shared(campus)
    with pytest.raises(PermissionDenied):
        campus.upload_document(shared["course_id"], "student_1", "student", "x.md", "text/markdown", b"content")
    with pytest.raises(PermissionDenied):
        campus.require_access(shared["course_id"], "student_2", "student")


def test_upload_security_duplicate_and_retrieval_isolation(campus):
    first = campus.create_course("课程一", "personal_course", "student_1", "student")
    second = campus.create_course("课程二", "personal_course", "student_1", "student")
    content = "# 独有概念\n蓝鲸算法只存在于课程一。".encode()
    campus.upload_document(first["course_id"], "student_1", "student", "safe.md", "text/markdown", content)
    with pytest.raises(ValidationError):
        campus.upload_document(first["course_id"], "student_1", "student", "again.md", "text/markdown", content)
    with pytest.raises(ValidationError):
        campus.upload_document(first["course_id"], "student_1", "student", "bad.pdf", "text/plain", b"not pdf")
    isolated = campus.ask(second["course_id"], "student_1", "student", "什么是蓝鲸算法？")
    assert isolated["refused"] is True
    found = campus.ask(first["course_id"], "student_1", "student", "什么是蓝鲸算法？")
    assert found["refused"] is False
    assert found["sources"][0]["source_file"] == "safe.md"


def test_quiz_profile_and_anonymous_class_analysis(campus):
    shared = add_shared(campus)
    qa = campus.ask(shared["course_id"], "student_1", "student", "监督学习使用什么样本？")
    quiz = campus.generate_quiz(shared["course_id"], "student_1", "student", qa["question_id"])
    answers = [item["answer"] for item in quiz["items"]]
    grade = campus.submit_quiz(shared["course_id"], "student_1", "student", qa["question_id"], quiz["items"], answers)
    assert grade["score"] == 100
    assert campus.profile(shared["course_id"], "student_1", "student")["attempts"]
    analysis = campus.class_analysis(shared["course_id"], "teacher_1")
    assert analysis["question_count"] == 1
    assert "student_1" not in str(analysis)


def test_virtual_course_upsert_owner_and_agent_contract(campus):
    first = campus.upsert_virtual_course("virtual_test", "第一版", "teacher_1")
    updated = campus.upsert_virtual_course("virtual_test", "第二版", "teacher_1")
    assert first["course_id"] == updated["course_id"] and updated["course_name"] == "第二版"
    with pytest.raises(PermissionDenied):
        campus.upsert_virtual_course("virtual_test", "越权修改", "teacher_2")

    agent = CampusAgentService(campus)
    response = agent.invoke({"request_id":"r1", "agent":"student_assistant", "action":"personal_course_create",
                             "actor":{"user_id":"s1","role":"student"}, "input":{"course_name":"Agent 课程"}})
    assert response.status == "success"
    unknown = agent.invoke({"request_id":"r2", "agent":"student_assistant", "action":"unknown_action",
                            "actor":{"user_id":"s1","role":"student"}})
    assert unknown.to_dict()["status"] == "not_implemented"


def test_agent_controlled_upload(campus):
    agent = CampusAgentService(campus)
    created = agent.invoke({"request_id":"r1", "agent":"student_assistant", "action":"personal_course_create",
                            "actor":{"user_id":"s1","role":"student"}, "input":{"course_name":"上传课"}}).data
    payload = base64.b64encode("# 内容\n这是受控上传。".encode()).decode()
    uploaded = agent.invoke({"request_id":"r2", "agent":"student_assistant", "action":"student_document_upload",
                             "actor":{"user_id":"s1","role":"student"}, "scope":{"course_id":created["course_id"]},
                             "input":{"file_name":"a.md","mime_type":"text/markdown","content_base64":payload}})
    assert uploaded.status == "success" and uploaded.data["chunk_count"] == 1


def test_report_exports(campus):
    shared = add_shared(campus)
    assert campus.export_class_csv(shared["course_id"], "teacher_1").startswith(b"\xef\xbb\xbf")
    assert campus.export_class_excel(shared["course_id"], "teacher_1").startswith(b"PK")
    assert campus.export_class_word(shared["course_id"], "teacher_1").startswith(b"PK")


def test_agent_uploads_and_extracts_pdf_docx_pptx(campus):
    from docx import Document
    from pptx import Presentation
    from pypdf import PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    agent = CampusAgentService(campus)
    course = campus.create_course("多格式课程", "personal_course", "student_1", "student")

    pdf_output = io.BytesIO()
    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    font = DictionaryObject({NameObject("/Type"): NameObject("/Font"),
                             NameObject("/Subtype"): NameObject("/Type1"),
                             NameObject("/BaseFont"): NameObject("/Helvetica")})
    resources = DictionaryObject({NameObject("/Font"): DictionaryObject({NameObject("/F1"): writer._add_object(font)})})
    content = DecodedStreamObject(); content.set_data(b"BT /F1 12 Tf 72 720 Td (PDF backend text) Tj ET")
    page[NameObject("/Resources")] = resources
    page[NameObject("/Contents")] = writer._add_object(content)
    writer.write(pdf_output)

    docx_output = io.BytesIO(); document = Document(); document.add_heading("Word section", 1)
    document.add_paragraph("Word backend text"); document.save(docx_output)

    pptx_output = io.BytesIO(); presentation = Presentation(); slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PPT section"; slide.placeholders[1].text = "PPT backend text"; presentation.save(pptx_output)

    files = [
        ("lesson.pdf", "application/pdf", pdf_output.getvalue()),
        ("lesson.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", docx_output.getvalue()),
        ("lesson.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", pptx_output.getvalue()),
    ]
    for index, (name, mime, content_bytes) in enumerate(files):
        response = agent.invoke({
            "request_id": f"format_{index}", "agent": "student_assistant", "action": "student_document_upload",
            "actor": {"user_id": "student_1", "role": "student"}, "scope": {"course_id": course["course_id"]},
            "input": {"file_name": name, "mime_type": mime,
                      "content_base64": base64.b64encode(content_bytes).decode("ascii")},
        })
        assert response.status == "success", response.message
        assert response.data["chunk_count"] >= 1

    documents = campus.list_documents(course["course_id"], "student_1", "student")
    assert len(documents) == 3
    assert all(row["text_preview"] for row in documents)


def test_qwen_provider_uses_real_compatible_endpoint_contract():
    class FakeResponse:
        status_code = 200
        text = ""
        def json(self): return {"choices": [{"message": {"content": "真实模型响应"}}]}

    provider = QwenProvider("sk-server-only", "https://ws-c4qflt1k6x8xwd4f.cn-beijing.maas.aliyuncs.com/compatible-mode/v1", "qwen-plus")
    with patch.object(provider.session, "post", return_value=FakeResponse()) as mocked:
        assert provider.generate("system", "user") == "真实模型响应"
    assert mocked.call_args.args[0] == "https://ws-c4qflt1k6x8xwd4f.cn-beijing.maas.aliyuncs.com/compatible-mode/v1/chat/completions"
    assert mocked.call_args.kwargs["json"]["model"] == "qwen-plus"
    assert mocked.call_args.kwargs["headers"]["Authorization"] == "Bearer sk-server-only"


def test_student_memory_minimum_loop_and_teacher_disabled(campus):
    agent = CampusAgentService(campus)
    course = campus.create_course("背诵课程", "personal_course", "student_1", "student")
    uploaded = campus.upload_document(course["course_id"], "student_1", "student", "material.txt", "text/plain",
                                      "监督学习使用带标签样本训练模型，并根据误差更新参数。".encode())
    built = agent.invoke({"request_id":"m1","agent":"student_assistant","action":"knowledge_blocks_build",
                          "actor":{"user_id":"student_1","role":"student"},"scope":{"course_id":course["course_id"]},
                          "input":{"document_id":uploaded["document_id"]}})
    assert built.status == "success" and built.data[0]["keywords"]
    block_id = built.data[0]["block_id"]
    cloze = agent.invoke({"request_id":"m2","agent":"student_assistant","action":"cloze_generate",
                          "actor":{"user_id":"student_1","role":"student"},"scope":{"course_id":course["course_id"]},
                          "input":{"block_id":block_id}})
    assert cloze.data["blank_count"] == 2
    cloze_grade = agent.invoke({"request_id":"m2b","agent":"student_assistant","action":"cloze_submit",
                                "actor":{"user_id":"student_1","role":"student"},"scope":{"course_id":course["course_id"]},
                                "input":{"block_id":block_id,"extra_keywords":[],"responses":["监督学习","错误"]}})
    assert cloze_grade.status == "success" and cloze_grade.data["score"] == 50
    assert cloze_grade.data["error_points"]
    evaluated = agent.invoke({"request_id":"m3","agent":"student_assistant","action":"recitation_evaluate",
                               "actor":{"user_id":"student_1","role":"student"},"scope":{"course_id":course["course_id"]},
                               "input":{"block_id":block_id,"recited_text":"监督学习使用带标签样本训练模型。"}})
    assert evaluated.data["score"] == 80
    questions = agent.invoke({"request_id":"m4","agent":"student_assistant","action":"memory_questions_generate",
                               "actor":{"user_id":"student_1","role":"student"},"scope":{"course_id":course["course_id"]},
                               "input":{"count":3}})
    assert questions.status == "success" and questions.data[0]["answer"] == "带标签样本"
    graded = agent.invoke({"request_id":"m5","agent":"student_assistant","action":"memory_questions_submit",
                           "actor":{"user_id":"student_1","role":"student"},"scope":{"course_id":course["course_id"]},
                           "input":{"questions":questions.data,"responses":["带标签样本"]}})
    assert graded.status == "success" and graded.data["score"] == 100
    wrong_grade = agent.invoke({"request_id":"m5b","agent":"student_assistant","action":"memory_questions_submit",
                                 "actor":{"user_id":"student_1","role":"student"},"scope":{"course_id":course["course_id"]},
                                 "input":{"questions":questions.data,"responses":["错误答案"]}})
    assert wrong_grade.status == "success" and wrong_grade.data["score"] == 0
    dashboard = agent.invoke({"request_id":"m6","agent":"student_assistant","action":"student_dashboard",
                              "actor":{"user_id":"student_1","role":"student"},"scope":{"course_id":course["course_id"]}})
    assert dashboard.data["recitation_book"] and dashboard.data["practice_average"] == 50
    assert dashboard.data["wrong_question_book"][0]["student_answer"] == "错误答案"

    export_actions = {
        "memory_workbook_export":{"course_name":"背诵课程","questions":questions.data},
        "recitation_book_export":{"course_name":"背诵课程"},
        "wrong_question_book_export":{"course_name":"背诵课程"},
    }
    for index, (action, input_data) in enumerate(export_actions.items()):
        exported = agent.invoke({"request_id":f"export_{index}","agent":"student_assistant","action":action,
                                 "actor":{"user_id":"student_1","role":"student"},
                                 "scope":{"course_id":course["course_id"]},"input":input_data})
        assert exported.status == "success"
        docx_bytes = base64.b64decode(exported.data["content_base64"])
        with zipfile.ZipFile(io.BytesIO(docx_bytes)) as archive:
            document_xml = archive.read("word/document.xml").decode("utf-8")
            styles_xml = archive.read("word/styles.xml").decode("utf-8")
        assert "宋体" in document_xml
        assert "宋体" in styles_xml
    teacher_course = campus.create_course("教师共享课程", "shared_course", "teacher_1", "teacher")
    enabled = agent.invoke({"request_id":"t1","agent":"teacher_assistant","action":"teaching_report",
                            "actor":{"user_id":"teacher_1","role":"teacher"},
                            "scope":{"course_id":teacher_course["course_id"]}})
    assert enabled.status == "success"

    deleted = agent.invoke({"request_id":"m7","agent":"student_assistant","action":"personal_course_delete",
                            "actor":{"user_id":"student_1","role":"student"},"scope":{"course_id":course["course_id"]}})
    assert deleted.status == "success"
    with pytest.raises(NotFound):
        campus.get_course(course["course_id"])


def test_agent_returns_model_network_error_without_crashing(tmp_path):
    class BrokenProvider(LLMProvider):
        def generate(self, system_prompt: str, user_prompt: str) -> str:
            raise RuntimeError("TLS handshake timeout")

    service = CampusService(LearningDatabase(tmp_path / "broken.db"), tmp_path / "uploads", BrokenProvider)
    course = service.create_course("网络异常课程", "personal_course", "student_1", "student")
    uploaded = service.upload_document(course["course_id"], "student_1", "student", "a.txt", "text/plain",
                                       "可供分块的学习材料内容，包含足够的文字。".encode())
    response = CampusAgentService(service).invoke({
        "request_id":"broken","agent":"student_assistant","action":"knowledge_blocks_build",
        "actor":{"user_id":"student_1","role":"student"},"scope":{"course_id":course["course_id"]},
        "input":{"document_id":uploaded["document_id"]},
    })
    assert response.status == "error"
    assert "TLS handshake timeout" in response.message


def test_import_question_bank_with_and_without_answers(campus):
    from openpyxl import Workbook

    agent = CampusAgentService(campus)
    course = campus.create_course("导入题库课程", "personal_course", "student_1", "student")
    text_payload = "1. 监督学习使用什么样本？\n答案：带标签样本\n2. 说明监督学习的含义。"
    imported = agent.invoke({
        "request_id":"bank_txt", "agent":"student_assistant", "action":"question_bank_import",
        "actor":{"user_id":"student_1","role":"student"}, "scope":{"course_id":course["course_id"]},
        "input":{"file_name":"题库.txt","mime_type":"text/plain",
                 "content_base64":base64.b64encode(text_payload.encode()).decode("ascii")},
    })
    assert imported.status == "success" and len(imported.data) == 2
    assert imported.data[0]["answer_source"] == "imported"
    assert imported.data[1]["answer"] == "" and imported.data[1]["answer_source"] == "ai_judge"

    workbook = Workbook(); sheet = workbook.active
    sheet.append(["题目", "选项", "答案"]); sheet.append(["监督学习使用什么样本？", "带标签样本|无标签样本", "带标签样本"])
    output = io.BytesIO(); workbook.save(output)
    excel_import = agent.invoke({
        "request_id":"bank_xlsx", "agent":"student_assistant", "action":"question_bank_import",
        "actor":{"user_id":"student_1","role":"student"}, "scope":{"course_id":course["course_id"]},
        "input":{"file_name":"题库.xlsx","mime_type":"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                 "content_base64":base64.b64encode(output.getvalue()).decode("ascii")},
    })
    assert excel_import.status == "success" and excel_import.data[0]["source_file"] == "题库.xlsx"


def test_question_bank_retries_when_model_json_is_truncated(tmp_path):
    class TruncatingQuestionBankProvider(LLMProvider):
        calls = 0

        def generate(self, system_prompt: str, user_prompt: str) -> str:
            self.__class__.calls += 1
            if "题库结构化整理智能体" not in system_prompt:
                return "{}"
            source = user_prompt.split("完整题目：\n", 1)[-1]
            if len(source) > 900:
                return (
                    '[{"type":"判断题","question":"机器学习是人工智能的一个子领域。",'
                    '"options":[],"answer":"✓","explanation":"","knowledge_point":"人工智能"},'
                    '{"type":"判断题","question":"监督学习需要标注数据。","options":[]'
                )
            return json.dumps([{
                "type":"判断题", "question":f"片段题目 {source[:8]}", "options":[],
                "answer":"✓", "explanation":"", "knowledge_point":"测试",
            }], ensure_ascii=False)

    service = CampusService(
        LearningDatabase(tmp_path / "truncated.db"), tmp_path / "uploads",
        TruncatingQuestionBankProvider,
    )
    course = service.create_course("长题库", "personal_course", "student_1", "student")
    long_bank = "\n".join(f"{index}. 这是一道用于测试自动分段重试的判断题。" for index in range(100))
    response = CampusAgentService(service).invoke({
        "request_id":"truncated_bank", "agent":"student_assistant", "action":"question_bank_import",
        "actor":{"user_id":"student_1","role":"student"}, "scope":{"course_id":course["course_id"]},
        "input":{"file_name":"长题库.txt","mime_type":"text/plain",
                 "content_base64":base64.b64encode(long_bank.encode()).decode("ascii")},
    })
    assert response.status == "success", response.message
    assert response.data
    assert all(item["options"] == ["正确", "错误"] for item in response.data)
    assert all(item["answer"] == "正确" for item in response.data)
    assert TruncatingQuestionBankProvider.calls >= 3


@pytest.mark.parametrize("raw", [
    "对", "正确", "√", "✓", "✔", "T", "t", "True", "TRUE", "是", "真", "Y", "Yes", 1, True,
])
def test_judgment_true_answer_variants(raw):
    assert _normalize_judgment_answer(raw) == "正确"


@pytest.mark.parametrize("raw", [
    "错", "错误", "×", "✕", "✖", "✗", "F", "f", "False", "FALSE", "否", "假", "N", "No", 0, False,
])
def test_judgment_false_answer_variants(raw):
    assert _normalize_judgment_answer(raw) == "错误"


def test_judgment_question_gets_options_and_extracts_inline_answer():
    question = _normalize_question_item({
        "type":"判断题", "question":"监督学习需要标注数据。（√）",
        "options":[], "answer":"", "knowledge_point":"监督学习",
    })
    assert question["question"] == "监督学习需要标注数据。"
    assert question["options"] == ["正确", "错误"]
    assert question["answer"] == "正确"
