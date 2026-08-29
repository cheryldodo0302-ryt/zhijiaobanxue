"""Local-first structure and knowledge-boundary utilities.

The ingestion pipeline deliberately keeps parser output immutable.  This module
only adds metadata and candidate records around source blocks; it never asks an
LLM to rewrite document content.
"""

from __future__ import annotations

import hashlib
import json
import re
from pathlib import Path
from typing import Any


REGION_TYPES = {
    "knowledge", "example", "case", "exercise", "solution",
    "introduction", "summary", "reference", "other",
}

_REGION_PATTERNS: tuple[tuple[str, tuple[str, ...]], ...] = (
    ("exercise", ("习题", "练习", "作业", "思考题", "课后题", "exercise", "quiz")),
    ("solution", ("解答", "答案", "解析", "参考答案", "solution", "answer")),
    ("example", ("例题", "例子", "示例", "例：", "例:", "example")),
    ("case", ("案例", "实例", "案例分析", "case study", "case")),
    ("introduction", ("引言", "绪论", "导论", "概述", "introduction")),
    ("summary", ("小结", "总结", "本章总结", "summary", "conclusion")),
    ("reference", ("参考文献", "参考资料", "references", "bibliography")),
)

_KNOWLEDGE_TYPE_PATTERNS: tuple[tuple[str, tuple[str, ...]], ...] = (
    ("definition", ("定义", "概念", "术语", "definition")),
    ("theorem", ("定理", "theorem")),
    ("property", ("性质", "property")),
    ("principle", ("原理", "定律", "原则", "principle")),
    ("formula", ("公式", "方程", "公式推导", "formula", "equation")),
    ("rule", ("规则", "法则", "rule")),
    ("procedure", ("步骤", "流程", "procedure")),
    ("method", ("方法", "算法", "method", "algorithm")),
    ("classification", ("分类", "类型", "类别", "classification")),
    ("comparison", ("比较", "区别", "对比", "comparison")),
    ("table", ("表格", "对照表", "table")),
    ("fact", ("事实", "事实数据", "fact")),
)


def block_text(block: dict[str, Any]) -> str:
    return str(
        block.get("plain_text")
        or block.get("markdown")
        or block.get("latex")
        or ""
    ).strip()


def block_raw(block: dict[str, Any]) -> dict[str, Any]:
    """Return parser metadata from either an in-memory or persisted block."""
    value = block.get("raw")
    if isinstance(value, dict):
        return value
    try:
        decoded = json.loads(block.get("raw_payload_json") or "{}")
    except (json.JSONDecodeError, TypeError):
        return {}
    return decoded if isinstance(decoded, dict) else {}


def normalize_title(value: str) -> str:
    value = re.sub(r"^\s*#+\s*", "", str(value or "")).strip()
    value = re.sub(
        r"^\s*(?:第\s*[0-9一二三四五六七八九十百千万]+\s*[章节篇部]|[0-9]+(?:\.[0-9]+)*[.)、]?)\s*",
        "",
        value,
        flags=re.I,
    )
    return re.sub(r"\s+", " ", value).strip(" ：:.-")


def _is_heading(block: dict[str, Any]) -> bool:
    raw = block_raw(block)
    if str(raw.get("source_kind") or "").lower() == "pptx":
        # Bullets such as "3. ..." are common inside slide bodies.  For PPTX,
        # only a shape that Fast Inspect identified as the title bar may open a
        # new structural group.
        return bool(raw.get("is_slide_title_block"))
    if str(block.get("block_type") or "") == "title":
        return True
    text = block_text(block)
    return bool(re.match(r"^\s*#{1,6}\s+", str(block.get("markdown") or ""))) or bool(
        re.match(
            r"^\s*(?:第\s*[0-9一二三四五六七八九十百千万]+\s*[章节篇部]|[0-9]+(?:\.[0-9]+)*[.)、])\s*\S+",
            text,
            flags=re.I,
        )
    )


def _heading_level(text: str, block: dict[str, Any] | None = None) -> int:
    raw = block_raw(block or {})
    if str(raw.get("source_kind") or "").lower() == "pptx":
        try:
            hint = int(raw.get("ppt_title_level_hint") or 0)
        except (TypeError, ValueError):
            hint = 0
        if hint:
            return max(1, min(3, hint))
    if re.match(r"^\s*#{1,6}\s+", text):
        return min(6, len(text) - len(text.lstrip("#")))
    if re.match(r"^\s*第\s*[0-9一二三四五六七八九十百千万]+\s*章", text):
        return 1
    if re.match(r"^\s*第\s*[0-9一二三四五六七八九十百千万]+\s*节", text):
        return 2
    match = re.match(r"^\s*([0-9]+(?:\.[0-9]+)*)[.)、]", text)
    if match:
        return min(6, match.group(1).count(".") + 1)
    return 3


def _is_structural_heading(block: dict[str, Any]) -> bool:
    markdown = str(block.get("markdown") or "").strip()
    raw = block_raw(block)
    try:
        raw_heading_level = int(raw.get("heading_level") or 0)
    except (AttributeError, TypeError, ValueError):
        raw_heading_level = 0
    return (
        str(block.get("block_type") or "") == "title"
        or bool(re.match(r"^##(?!#)\s+", markdown))
        or raw_heading_level == 2
    )


def _numbered_secondary_title(block: dict[str, Any]) -> str:
    """Return an ``x.x`` candidate, excluding deeper ``x.x.x`` numbers."""
    if not _is_heading(block):
        return ""
    text = re.sub(r"^\s*#{1,6}\s*", "", block_text(block)).strip()
    if not text or len(text) > 180:
        return ""
    if not re.match(r"^\d+\.\d+(?!\.\d)\s*\S", text):
        return ""
    return re.sub(r"\s+", " ", text)[:180]


class RegionClassifier:
    """Classify regions using headings and parent inheritance.

    A heading starts a region and following blocks inherit it until another
    heading starts a new region.  This is intentionally deterministic so a
    teacher can understand why a block was routed out of the knowledge layer.
    """

    def classify(self, blocks: list[dict[str, Any]]) -> list[dict[str, Any]]:
        active_region = "knowledge"
        active_parent = ""
        active_reason = "默认知识正文"
        active_confidence = 0.45
        result: list[dict[str, Any]] = []
        for block in blocks:
            item = dict(block)
            text = block_text(item)
            page_type = str(item.get("page_type") or "CONTENT").upper()
            explicit = self._region_hint(text) if _is_heading(item) else None
            if explicit:
                active_region = explicit[0]
                active_reason = explicit[1]
                active_confidence = 0.95
                active_parent = str(item.get("block_id") or "")
            elif _is_heading(item):
                active_region = "knowledge"
                active_reason = "未命中特殊区域标题，作为知识章节"
                active_confidence = 0.7
                active_parent = str(item.get("block_id") or "")

            if page_type in {"COVER", "BACK_COVER", "TOC", "ENDING", "BLANK", "VISUAL"}:
                region = "other"
                reason = f"页面路由为 {page_type}"
                confidence = 0.98
                include = False
            else:
                region = active_region if active_region in REGION_TYPES else "other"
                reason = active_reason
                confidence = active_confidence
                include = region not in {"example", "case", "exercise", "solution", "reference", "other"}
                if bool(item.get("include_as_knowledge", True)) is False:
                    include = False
            destination = "unclassified" if include else "excluded"
            if region in {"example", "case", "exercise", "solution"}:
                destination = "question_bank"
            item.update({
                "region_type": region,
                "region_confidence": confidence,
                "region_reason": reason,
                "parent_region_block_id": active_parent if active_parent != str(item.get("block_id") or "") else "",
                "include_as_knowledge": bool(include),
                "content_destination": destination,
            })
            result.append(item)
        return result

    @staticmethod
    def _region_hint(text: str) -> tuple[str, str] | None:
        lowered = text.lower()
        for region, terms in _REGION_PATTERNS:
            if any(term.lower() in lowered for term in terms):
                return region, f"标题命中 {region} 区域规则"
        return None


class StructureBuilder:
    """Build a lightweight outline and report TOC/heading mismatches."""

    def build(self, blocks: list[dict[str, Any]], pages: list[dict[str, Any]] | None = None) -> dict[str, Any]:
        outline: list[dict[str, Any]] = []
        stack: list[dict[str, Any]] = []
        chapter_paths: dict[str, list[str]] = {}
        for block in blocks:
            if str(block.get("page_type") or "").upper() in {"TOC", "COVER", "BACK_COVER", "ENDING", "BLANK", "VISUAL"}:
                continue
            if str(block.get("region_type") or "knowledge") in {"other", "reference", "example", "case", "exercise", "solution"}:
                continue
            text = block_text(block)
            if _is_heading(block):
                title = normalize_title(text)
                if not title:
                    continue
                level = _heading_level(text, block)
                while stack and int(stack[-1]["level"]) >= level:
                    stack.pop()
                node = {
                    "node_id": f"outline_{len(outline) + 1:04d}",
                    "level": level,
                    "title": title,
                    "page_start": int(block.get("page_number") or 1),
                    "page_end": int(block.get("page_number") or 1),
                    "block_ids": [str(block.get("block_id") or "")],
                    "children": [],
                }
                if stack:
                    stack[-1]["children"].append(node)
                else:
                    outline.append(node)
                stack.append(node)
                chapter_paths[str(block.get("block_id") or "")] = [x["title"] for x in stack]
                continue
            if stack:
                stack[-1]["block_ids"].append(str(block.get("block_id") or ""))
                stack[-1]["page_end"] = max(stack[-1]["page_end"], int(block.get("page_number") or 1))
                chapter_paths[str(block.get("block_id") or "")] = [x["title"] for x in stack]

        toc_entries = self._toc_entries(blocks, pages or [])
        heading_titles = {normalize_title(node["title"]).lower() for node in self._flatten(outline)}
        warnings: list[dict[str, Any]] = []
        for entry in toc_entries:
            if normalize_title(entry).lower() not in heading_titles:
                warnings.append({
                    "code": "STRUCTURE_WARNING",
                    "type": "toc_heading_missing",
                    "expected_title": normalize_title(entry),
                    "message": f"目录条目未在正文标题中匹配：{normalize_title(entry)}",
                })
        if toc_entries and not outline:
            warnings.append({
                "code": "STRUCTURE_WARNING",
                "type": "toc_without_headings",
                "expected_title": "",
                "message": "检测到目录，但正文未解析出可用章节标题",
            })
        return {
            "status": "warning" if warnings else "ok",
            "outline": outline,
            "toc_entries": toc_entries,
            "warnings": warnings,
            "block_chapter_paths": chapter_paths,
        }

    @staticmethod
    def _flatten(nodes: list[dict[str, Any]]) -> list[dict[str, Any]]:
        result: list[dict[str, Any]] = []
        for node in nodes:
            result.append(node)
            result.extend(StructureBuilder._flatten(node.get("children") or []))
        return result

    @staticmethod
    def _toc_entries(blocks: list[dict[str, Any]], pages: list[dict[str, Any]]) -> list[str]:
        entries: list[str] = []
        page_types = {int(page.get("page_number") or 0): str(page.get("page_type") or "").upper() for page in pages}
        for block in blocks:
            page_type = str(block.get("page_type") or page_types.get(int(block.get("page_number") or 0), "")).upper()
            if page_type != "TOC":
                continue
            for line in block_text(block).splitlines():
                line = re.sub(r"\.{2,}\s*\d+\s*$", "", line).strip()
                if re.match(r"^(?:第\s*\S+\s*[章节]|\d+(?:\.\d+)*[.)、])", line):
                    entries.append(normalize_title(line))
        return list(dict.fromkeys(x for x in entries if x))


class KnowledgeBoundaryExtractor:
    """Create source-block-first knowledge candidates."""

    def extract(self, blocks: list[dict[str, Any]], document_id: str = "") -> list[dict[str, Any]]:
        ordered = sorted(blocks, key=lambda row: (int(row.get("page_index", int(row.get("page_number") or 1) - 1)), int(row.get("block_order") or 0)))
        presentation_groups = self.extract_presentation_title_groups(ordered, document_id)
        if presentation_groups:
            return presentation_groups
        numbered_sections = self.extract_numbered_sections(ordered, document_id)
        if numbered_sections:
            return numbered_sections
        candidates: list[dict[str, Any]] = []
        current: list[dict[str, Any]] = []
        for block in ordered:
            if not bool(block.get("include_as_knowledge", True)) or str(block.get("region_type") or "knowledge") != "knowledge":
                if current:
                    candidates.append(self._candidate(current, document_id))
                    current = []
                continue
            if _is_heading(block) and self._is_container_heading(block_text(block)):
                if current:
                    candidates.append(self._candidate(current, document_id))
                    current = []
                continue
            if _is_heading(block) and current:
                candidates.append(self._candidate(current, document_id))
                current = []
            current.append(block)
        if current:
            candidates.append(self._candidate(current, document_id))
        return [candidate for candidate in candidates if candidate["markdown_content"].strip()]

    def extract_presentation_title_groups(
        self, blocks: list[dict[str, Any]], document_id: str = ""
    ) -> list[dict[str, Any]]:
        """Group PPT source blocks by consecutive, Fast-Inspect title bars.

        A slide-body bullet that happens to start with a number is never treated
        as a boundary here.  Repeated titles on consecutive slides remain one
        review candidate, while a title that reappears later starts a new group
        so the candidate queue always follows presentation order.
        """
        ordered = sorted(
            blocks,
            key=lambda row: (
                int(row.get("page_index", int(row.get("page_number") or 1) - 1)),
                int(row.get("block_order") or 0),
            ),
        )
        if not any(
            str(block_raw(block).get("source_kind") or "").lower() == "pptx"
            and bool(block_raw(block).get("is_slide_title_block"))
            for block in ordered
        ):
            return []

        candidates: list[dict[str, Any]] = []
        current: list[dict[str, Any]] = []
        current_key = ""

        def flush() -> None:
            nonlocal current
            if current:
                candidates.append(
                    self._candidate(current, document_id, preserve_heading_title=True)
                )
                current = []

        for block in ordered:
            raw = block_raw(block)
            is_title = (
                str(raw.get("source_kind") or "").lower() == "pptx"
                and bool(raw.get("is_slide_title_block"))
            )
            region = str(block.get("region_type") or "knowledge")
            eligible = bool(block.get("include_as_knowledge", True)) and region not in {
                "example", "case", "exercise", "solution", "reference", "other",
            }
            if is_title:
                title = str(raw.get("ppt_slide_title") or block_text(block)).strip()
                title = re.sub(r"^\s*#{1,6}\s*", "", title).strip()
                title_key = re.sub(r"\s+", " ", title).casefold()
                if not eligible:
                    flush()
                    current_key = ""
                    continue
                if current and title_key != current_key:
                    flush()
                if not current:
                    current = [block]
                    current_key = title_key
                elif str(block.get("block_id") or "") not in {
                    str(value.get("block_id") or "") for value in current
                }:
                    current.append(block)
                continue
            if current and eligible:
                current.append(block)
            elif current and not eligible:
                flush()
                current_key = ""
        flush()
        return [candidate for candidate in candidates if candidate["markdown_content"].strip()]

    def extract_numbered_sections(
        self, blocks: list[dict[str, Any]], document_id: str = ""
    ) -> list[dict[str, Any]]:
        """Group one review candidate from each numbered secondary heading.

        An ``x.x`` title owns every eligible source block below it, including
        ``x.x.x`` descendants, until the next level-one or level-two heading.
        This keeps the teacher approval queue faithful to the textbook outline
        instead of producing one candidate for every paragraph or tertiary title.
        """
        ordered = sorted(
            blocks,
            key=lambda row: (
                int(row.get("page_index", int(row.get("page_number") or 1) - 1)),
                int(row.get("block_order") or 0),
            ),
        )
        if not any(_numbered_secondary_title(block) for block in ordered):
            return []

        candidates: list[dict[str, Any]] = []
        current: list[dict[str, Any]] = []
        last_section_number: tuple[int, int] | None = None

        def flush() -> None:
            nonlocal current
            if current:
                candidates.append(
                    self._candidate(current, document_id, preserve_heading_title=True)
                )
                current = []

        for block in ordered:
            section_title = _numbered_secondary_title(block)
            region = str(block.get("region_type") or "knowledge")
            eligible = bool(block.get("include_as_knowledge", True)) and region not in {
                "example", "case", "exercise", "solution", "reference", "other",
            }
            section_number_match = re.match(r"^(\d+)\.(\d+)(?!\.\d)", section_title)
            section_number = (
                (int(section_number_match.group(1)), int(section_number_match.group(2)))
                if section_number_match else None
            )
            is_forward = bool(section_number) and (
                last_section_number is None or section_number > last_section_number
            )
            is_expected_next = bool(section_number and last_section_number) and (
                (section_number[0] == last_section_number[0]
                 and section_number[1] == last_section_number[1] + 1)
                or (section_number[0] == last_section_number[0] + 1
                    and section_number[1] == 1)
            )
            reliable_boundary = _is_structural_heading(block) or is_expected_next
            if section_title and eligible and is_forward and reliable_boundary:
                flush()
                current = [block]
                last_section_number = section_number
                continue
            markdown = str(block.get("markdown") or "").strip()
            is_chapter = bool(re.match(r"^#(?!#)\s+", markdown)) or bool(re.match(
                r"^\s*第\s*[0-9一二三四五六七八九十百千万]+\s*[章篇部]", block_text(block)
            ))
            if _is_heading(block) and is_chapter:
                flush()
                continue
            if current and eligible:
                current.append(block)
        flush()
        return [candidate for candidate in candidates if candidate["markdown_content"].strip()]

    @staticmethod
    def _is_container_heading(text: str) -> bool:
        normalized = normalize_title(text)
        heading = re.sub(r"^\s*#+\s*", "", text).strip()
        return bool(re.match(r"^(?:第\s*[0-9一二三四五六七八九十百千万]+\s*[章节篇部]|\d+(?:\.\d+)?[.)、])", heading)) and not any(
            term in normalized for _, terms in _KNOWLEDGE_TYPE_PATTERNS for term in terms if len(term) > 1
        )

    def _candidate(
        self,
        source_blocks: list[dict[str, Any]],
        document_id: str,
        *,
        preserve_heading_title: bool = False,
    ) -> dict[str, Any]:
        ids = [str(block.get("block_id") or "") for block in source_blocks]
        seed = f"{document_id}|{'|'.join(ids)}"
        candidate_id = f"kc_{hashlib.sha256(seed.encode('utf-8')).hexdigest()[:24]}"
        first_raw = block_raw(source_blocks[0])
        source_title = (
            str(first_raw.get("ppt_slide_title") or block_text(source_blocks[0]))
            if str(first_raw.get("source_kind") or "").lower() == "pptx"
            else block_text(source_blocks[0])
        )
        title = (
            re.sub(r"^\s*#{1,6}\s*", "", source_title).strip()[:160]
            if preserve_heading_title
            else normalize_title(source_title)[:160]
        )
        if not title:
            title = block_text(source_blocks[0])[:80]
        joined = "\n\n".join(
            str(block.get("markdown") or block.get("latex") or block_text(block)).strip()
            for block in source_blocks
        )
        title_text = " ".join(block_text(block) for block in source_blocks[:2])
        knowledge_type = "concept"
        for kind, terms in _KNOWLEDGE_TYPE_PATTERNS:
            if any(term.lower() in title_text.lower() for term in terms):
                knowledge_type = kind
                break
        bboxes: list[Any] = []
        for block in source_blocks:
            bbox = block.get("bbox", [])
            if isinstance(bbox, str):
                try:
                    bbox = json.loads(bbox or "[]")
                except json.JSONDecodeError:
                    bbox = []
            if bbox:
                bboxes.append(bbox)
        pages = [int(block.get("page_number") or 1) for block in source_blocks]
        return {
            "candidate_id": candidate_id,
            "document_id": document_id,
            "title": title,
            "knowledge_type": knowledge_type,
            "source_block_ids": ids,
            "page_start": min(pages),
            "page_end": max(pages),
            "bbox": bboxes,
            "markdown_content": joined,
            "confidence": 0.9 if _is_heading(source_blocks[0]) else 0.65,
            "region_type": "knowledge",
            "chapter_path": self._chapter_path(source_blocks[0]),
        }

    @staticmethod
    def _chapter_path(block: dict[str, Any]) -> list[str]:
        value = block.get("chapter_path") or []
        if isinstance(value, str):
            try:
                value = json.loads(value or "[]")
            except json.JSONDecodeError:
                value = []
        return [str(item) for item in value] if isinstance(value, list) else []


class PptFastInspector:
    """Extract slide-level layout metadata without OCR or vision calls."""

    @classmethod
    def _leaf_shapes(cls, shapes: Any) -> list[Any]:
        """Flatten nested PowerPoint groups while retaining slide coordinates.

        Many templates place the visible top title inside two or more GROUP
        shapes.  ``slide.shapes`` only returns the outer group, whose ``text``
        is empty, so inspecting only the first level mistakes footer dates or
        body copy for the slide title.
        """
        leaves: list[Any] = []
        for shape in shapes:
            children = getattr(shape, "shapes", None)
            if children is not None:
                leaves.extend(cls._leaf_shapes(children))
            else:
                leaves.append(shape)
        return leaves

    def inspect(self, path: Path) -> list[dict[str, Any]]:
        from pptx import Presentation  # lazy: PDF/text tests do not need python-pptx

        presentation = Presentation(str(path))
        slides: list[dict[str, Any]] = []
        for index, slide in enumerate(presentation.slides):
            top_level_shapes = list(slide.shapes)
            shapes = self._leaf_shapes(top_level_shapes)
            ordered = sorted(shapes, key=lambda shape: (int(getattr(shape, "top", 0)), int(getattr(shape, "left", 0))))
            texts: list[dict[str, Any]] = []
            pictures = 0
            font_sizes: list[float] = []
            explicit_title = getattr(getattr(slide, "shapes", None), "title", None)
            explicit_title_id = int(getattr(explicit_title, "shape_id", -1)) if explicit_title is not None else -1
            for order, shape in enumerate(ordered):
                text = re.sub(r"\s+", " ", str(getattr(shape, "text", "") or "")).strip()
                shape_type = str(getattr(shape, "shape_type", ""))
                try:
                    shape_type_number = int(getattr(shape, "shape_type", 0) or 0)
                except (TypeError, ValueError):
                    shape_type_number = 0
                if "PICTURE" in shape_type.upper() or shape_type_number == 13:
                    pictures += 1
                shape_font_sizes: list[float] = []
                for paragraph in getattr(getattr(shape, "text_frame", None), "paragraphs", []) or []:
                    for run in getattr(paragraph, "runs", []) or []:
                        size = getattr(getattr(run, "font", None), "size", None)
                        if size is not None and getattr(size, "pt", None):
                            point_size = float(size.pt)
                            font_sizes.append(point_size)
                            shape_font_sizes.append(point_size)
                placeholder_type = ""
                try:
                    if bool(getattr(shape, "is_placeholder", False)):
                        placeholder_type = str(shape.placeholder_format.type)
                except (AttributeError, ValueError):
                    placeholder_type = ""
                shape_id = int(getattr(shape, "shape_id", order))
                texts.append({
                    "shape_id": shape_id,
                    "order": order,
                    "text": text,
                    "shape_type": shape_type,
                    "placeholder_type": placeholder_type,
                    "font_size_max": max(shape_font_sizes) if shape_font_sizes else None,
                    "is_title": shape_id == explicit_title_id,
                    "is_boilerplate": self._is_boilerplate(text),
                    "bbox_emu": [
                        int(getattr(shape, "left", 0)), int(getattr(shape, "top", 0)),
                        int(getattr(shape, "width", 0)), int(getattr(shape, "height", 0)),
                    ],
                })
            # Spatial order is more reliable than the TITLE placeholder in
            # templated decks: a lower body subheading is often the placeholder
            # while the visible top bar is a normal textbox inside a group.
            title_item = self._section_title_shape(texts)
            if title_item is None:
                title_item = self._top_banner_title_shape(
                texts, int(presentation.slide_width), int(presentation.slide_height)
                )
            if title_item is None:
                title_item = self._fallback_title_shape(
                    texts, int(presentation.slide_width), int(presentation.slide_height)
                )
                if title_item is not None:
                    title_item["is_title"] = True
            title = str((title_item or {}).get("text") or "").strip()
            text_count = sum(bool(item["text"]) for item in texts)
            complex_layout = len(shapes) > 8 or pictures > 0 or text_count > 5
            slide_type = self._slide_type(title, texts, index, len(presentation.slides))
            regions, layout_kind = self._reconstruct_regions(texts, int(presentation.slide_width), int(presentation.slide_height))
            reading_order = [
                item["shape_id"]
                for item in sorted(
                    (item for item in texts if item["text"]),
                    key=lambda item: (
                        0 if item.get("is_title") else 1,
                        item["bbox_emu"][1], item["bbox_emu"][0],
                    ),
                )
            ]
            slides.append({
                "slide_index": index,
                "slide_number": index + 1,
                "slide_type": slide_type,
                "parse_level": "NORMAL" if complex_layout else "FAST",
                "title": title[:300],
                "shape_count": len(shapes),
                "top_level_shape_count": len(top_level_shapes),
                "text_count": text_count,
                "picture_count": pictures,
                "font_size_min": min(font_sizes) if font_sizes else None,
                "font_size_max": max(font_sizes) if font_sizes else None,
                "reading_order": reading_order,
                "shapes": texts,
                "layout_kind": layout_kind,
                "regions": regions,
            })
        return slides

    @staticmethod
    def _reconstruct_regions(texts: list[dict[str, Any]], slide_width: int, slide_height: int) -> tuple[list[dict[str, Any]], str]:
        """Group shapes into coarse columns to preserve 2-D reading structure."""
        visible = [item for item in texts if item.get("text")]
        if not visible:
            return [], "empty"
        title_shapes = [item for item in visible if item.get("is_title")]
        content_shapes = [item for item in visible if not item.get("is_title")]
        if not content_shapes:
            return ([PptFastInspector._region_from_shapes("title", title_shapes, "title")] if title_shapes else []), "title_only"
        columns: dict[int, list[dict[str, Any]]] = {0: [], 1: [], 2: []}
        third = max(1, slide_width // 3)
        for shape in content_shapes:
            left, top, width, height = shape["bbox_emu"]
            center = left + max(0, width) // 2
            column = 0 if center < third else 1 if center < third * 2 else 2
            columns[column].append(shape)
        used = [items for items in columns.values() if items]
        if len(used) == 1:
            layout = "single_column"
        elif len(used) == 2:
            layout = "two_column"
        else:
            layout = "multi_column"
        regions: list[dict[str, Any]] = []
        if title_shapes:
            regions.append(PptFastInspector._region_from_shapes("title", title_shapes, "title"))
        for region_index, items in enumerate(used, 1):
            region = PptFastInspector._region_from_shapes(f"r{region_index:02d}", items, "content")
            _left, _top, width, height = region["bbox_emu"]
            region["relative_area"] = round(
                max(0, width) * max(0, height) / max(1, slide_width * slide_height), 6
            )
            regions.append(region)
        return regions, layout

    @staticmethod
    def _region_from_shapes(region_id: str, items: list[dict[str, Any]], role: str) -> dict[str, Any]:
        left = min(item["bbox_emu"][0] for item in items)
        top = min(item["bbox_emu"][1] for item in items)
        right = max(item["bbox_emu"][0] + item["bbox_emu"][2] for item in items)
        bottom = max(item["bbox_emu"][1] + item["bbox_emu"][3] for item in items)
        ordered = sorted(items, key=lambda item: (item["bbox_emu"][1], item["bbox_emu"][0]))
        return {
            "region_id": region_id, "role": role,
            "bbox_emu": [left, top, right - left, bottom - top],
            "shape_ids": [item["shape_id"] for item in ordered],
            "text": "\n".join(item["text"] for item in ordered if item["text"]),
            "relative_area": 0.0,
        }

    @staticmethod
    def _fallback_title_shape(
        texts: list[dict[str, Any]], slide_width: int, slide_height: int
    ) -> dict[str, Any] | None:
        """Choose a large section title when a slide has no top title bar."""
        candidates: list[tuple[float, dict[str, Any]]] = []
        for item in texts:
            text = str(item.get("text") or "").strip()
            if (not text or not PptFastInspector._looks_like_title(text)
                    or item.get("is_boilerplate") or PptFastInspector._is_numeric_marker(text)):
                continue
            left, top, width, _height = item["bbox_emu"]
            font_size = float(item.get("font_size_max") or 0)
            placeholder = str(item.get("placeholder_type") or "").upper()
            is_real_title_placeholder = "TITLE" in placeholder and "SUBTITLE" not in placeholder
            if font_size < 36 and not is_real_title_placeholder:
                continue
            score = 0.0
            if is_real_title_placeholder:
                score += 100.0
            score += max(0.0, 30.0 * (1.0 - top / max(1, slide_height)))
            score += min(20.0, 20.0 * width / max(1, slide_width))
            score += min(30.0, font_size / 2.0)
            if left < slide_width * 0.2:
                score += 5.0
            candidates.append((score, item))
        return max(candidates, key=lambda value: value[0])[1] if candidates else None

    @staticmethod
    def _top_banner_title_shape(
        texts: list[dict[str, Any]], slide_width: int, slide_height: int
    ) -> dict[str, Any] | None:
        candidates: list[tuple[tuple[float, float, float], dict[str, Any]]] = []
        for item in texts:
            text = str(item.get("text") or "").strip()
            if (not text or not PptFastInspector._looks_like_title(text)
                    or item.get("is_boilerplate") or PptFastInspector._is_numeric_marker(text)):
                continue
            left, top, width, _height = item["bbox_emu"]
            if top > slide_height * 0.17:
                continue
            # Tiny coordinates commonly belong to a diagram's nested local
            # coordinate system, not to a slide-level title textbox.
            if width < slide_width * 0.02:
                continue
            font_size = float(item.get("font_size_max") or 0)
            placeholder = str(item.get("placeholder_type") or "").upper()
            if font_size < 18 and "TITLE" not in placeholder:
                continue
            # Topmost meaningful text wins. Font size and width break ties.
            candidates.append((
                (float(top), -font_size, -float(width) / max(1, slide_width)), item
            ))
        return min(candidates, key=lambda value: value[0])[1] if candidates else None

    @staticmethod
    def _section_title_shape(texts: list[dict[str, Any]]) -> dict[str, Any] | None:
        """Recognise full-page dividers such as “学习内容” and “目录”."""
        candidates = [
            item for item in texts
            if float(item.get("font_size_max") or 0) >= 48
            and PptFastInspector._looks_like_title(str(item.get("text") or ""))
            and len(str(item.get("text") or "").strip()) <= 40
            and not item.get("is_boilerplate")
            and not PptFastInspector._is_numeric_marker(str(item.get("text") or ""))
        ]
        return max(candidates, key=lambda item: float(item.get("font_size_max") or 0)) if candidates else None

    @staticmethod
    def _is_numeric_marker(text: str) -> bool:
        return bool(re.fullmatch(r"[0-9０-９]+(?:[.、)]\s*)?", text.strip()))

    @staticmethod
    def _is_boilerplate(text: str) -> bool:
        lowered = re.sub(r"\s+", " ", str(text or "")).strip().lower()
        if not lowered:
            return False
        return (
            lowered.startswith("principle & application of dat")
            or bool(re.fullmatch(r"(?:january|february|march|april|may|june|july|august|september|october|november|december)\s+\d{1,2}(?:st|nd|rd|th)?", lowered))
            or lowered in {"温州医科大学", "wenzhou medical university"}
        )

    @staticmethod
    def _looks_like_title(text: str) -> bool:
        return len(text) <= 160 and "\n" not in text

    @staticmethod
    def _slide_type(title: str, texts: list[dict[str, Any]], index: int, total: int) -> str:
        lowered = title.lower()
        if index == 0 or any(term in lowered for term in ("目录", "contents", "封面", "title")):
            return "COVER"
        if any(term in lowered for term in ("谢谢", "感谢", "thank", "参考文献", "references")):
            return "ENDING"
        if len(texts) <= 2 and title:
            return "SECTION"
        return "COMPLEX_CONTENT" if len(texts) > 8 else "SIMPLE_CONTENT"
