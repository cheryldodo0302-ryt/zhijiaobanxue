from __future__ import annotations

import csv
import hashlib
import io
import json
import mimetypes
import re
import shutil
import uuid
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.decomposition import TruncatedSVD
from sklearn.preprocessing import normalize
from sklearn.metrics.pairwise import cosine_similarity

from config import DATA_DIR, MAX_EVIDENCE_CHARS, MAX_UPLOAD_BYTES, MIN_EVIDENCE_SCORE, TOP_K
from database import LearningDatabase
from llm_provider import LLMProvider, build_backend_provider
from security_utils import UnsafeUpload, validate_document_bytes
from skills.exercise import ExerciseItem, generate_exercises, grade_exercises
from skills.qa import answer_question, guide_question
from skills.retrieval import Evidence


COURSE_MATERIAL_LABELS = {
    "slides": "课件", "textbook": "教材", "syllabus": "教学大纲",
    "lesson_plan": "教案", "experiment": "实验资料", "question_bank": "题库",
    "knowledge_graph": "知识图谱", "teaching_schedule": "教学进度", "other": "其他",
}


class CampusError(Exception):
    pass


class PermissionDenied(CampusError):
    pass


class ValidationError(CampusError):
    pass


class NotFound(CampusError):
    pass


ALLOWED_FILES = {
    ".pdf": {"application/pdf"},
    ".docx": {"application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/octet-stream"},
    ".pptx": {"application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/octet-stream"},
    ".md": {"text/markdown", "text/plain", "application/octet-stream"},
    ".txt": {"text/plain", "application/octet-stream"},
}
def _json(value: Any) -> str:
    return json.dumps(value, ensure_ascii=False)


def _loads(value: str | None) -> Any:
    return json.loads(value or "[]")


def _safe_name(name: str) -> str:
    name = Path(name).name.strip().replace("\x00", "")
    stem = re.sub(r"[^\w\u4e00-\u9fff .()-]", "_", Path(name).stem)[:80].strip(" .") or "document"
    suffix = Path(name).suffix.lower()
    return f"{stem}{suffix}"


def _split_text(text: str, default_section: str = "正文", page_number: int | None = None) -> list[dict]:
    chunks: list[dict] = []
    section = default_section
    heading_level = 0
    heading_path: list[str] = []
    buffer: list[str] = []

    def flush(*, preserve_empty_heading: bool = False) -> None:
        content = "\n".join(buffer).strip()
        if not content:
            if preserve_empty_heading and heading_level:
                chunks.append({
                    "section": section,
                    "page_number": page_number,
                    "content": "",
                    "heading_level": heading_level,
                    "heading_path": list(heading_path),
                })
            return
        for start in range(0, len(content), 900):
            part = content[start:start + 1100].strip()
            if part:
                chunks.append({
                    "section": section,
                    "page_number": page_number,
                    "content": part,
                    "heading_level": heading_level,
                    "heading_path": list(heading_path),
                })

    for line in text.splitlines():
        heading = re.match(r"^(#{1,6})\s+(.+?)\s*$", line)
        if heading:
            flush(preserve_empty_heading=True)
            buffer.clear()
            heading_level = len(heading.group(1))
            section = heading.group(2).strip() or default_section
            heading_path = heading_path[:heading_level - 1]
            while len(heading_path) < heading_level - 1:
                heading_path.append("")
            heading_path.append(section)
        elif line.strip():
            buffer.append(line.strip())
    flush(preserve_empty_heading=True)
    return chunks


def parse_document(data: bytes, suffix: str) -> list[dict]:
    if suffix in {".md", ".txt"}:
        try:
            return _split_text(data.decode("utf-8-sig"))
        except UnicodeDecodeError as exc:
            raise ValidationError("文本文件必须使用 UTF-8 编码") from exc
    if suffix == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        if len(reader.pages) > 1000:
            raise ValidationError("PDF 页数超过 1000 页，请拆分后上传")
        chunks: list[dict] = []
        for index, page in enumerate(reader.pages, 1):
            chunks += _split_text(page.extract_text() or "", f"第 {index} 页", index)
        return chunks
    if suffix == ".docx":
        from docx import Document
        doc = Document(io.BytesIO(data))
        if len(doc.paragraphs) > 100000:
            raise ValidationError("DOCX 段落过多，请拆分后上传")
        lines = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                lines.append(("# " if paragraph.style.name.startswith("Heading") else "") + text)
        return _split_text("\n".join(lines))
    if suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        if len(prs.slides) > 1000:
            raise ValidationError("PPTX 页数超过 1000 页，请拆分后上传")
        chunks = []
        for index, slide in enumerate(prs.slides, 1):
            texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                chunks += _split_text("\n".join(texts), f"第 {index} 页", index)
        return chunks
    raise ValidationError("不支持的文件类型")


class ChunkRetriever:
    def __init__(self, rows: list[dict]):
        self.rows = rows
        self.chunks = rows
        self.vectorizer = TfidfVectorizer(analyzer="char", ngram_range=(1, 3), min_df=1, max_features=20000)
        self.matrix = self.vectorizer.fit_transform([row["content"] for row in rows]) if rows else None
        self.reducer = None
        self.faiss_index = None
        if self.matrix is not None and self.matrix.shape[0]:
            dimensions = min(256, self.matrix.shape[0] - 1, self.matrix.shape[1] - 1)
            dense = self.matrix.toarray()
            if dimensions >= 2:
                self.reducer = TruncatedSVD(n_components=dimensions, random_state=42)
                dense = self.reducer.fit_transform(self.matrix)
            dense = normalize(dense).astype("float32")
            try:
                import faiss

                self.faiss_index = faiss.IndexFlatIP(dense.shape[1])
                self.faiss_index.add(dense)
            except ImportError:
                self.faiss_index = None

    @staticmethod
    def _keywords(text: str) -> set[str]:
        runs = re.findall(r"[\u4e00-\u9fff]+", text)
        chinese = {run[i:i+n] for run in runs for n in (2, 3, 4) for i in range(max(len(run)-n+1, 0))}
        return chinese | set(re.findall(r"[A-Za-z0-9_-]{2,}", text.lower()))

    def search(self, query: str, top_k: int = 4) -> list[Evidence]:
        if not query.strip() or self.matrix is None:
            return []
        query_sparse = self.vectorizer.transform([query])
        scores = cosine_similarity(query_sparse, self.matrix)[0]
        candidate_indexes = list(range(len(self.rows)))
        if self.faiss_index is not None:
            query_dense = self.reducer.transform(query_sparse) if self.reducer is not None else query_sparse.toarray()
            query_dense = normalize(query_dense).astype("float32")
            candidate_count = min(len(self.rows), max(20, int(top_k) * 8))
            faiss_scores, faiss_indexes = self.faiss_index.search(query_dense, candidate_count)
            candidate_indexes = [int(index) for index in faiss_indexes[0] if index >= 0]
            scores = [0.0] * len(self.rows)
            for index, score in zip(candidate_indexes, faiss_scores[0]):
                scores[index] = float(score)
        query_keys = self._keywords(query)
        ranked = []
        for index in candidate_indexes:
            row = self.rows[index]
            keys = self._keywords(row["content"] + " " + row["section"])
            score = .72 * float(scores[index]) + .28 * len(query_keys & keys) / max(len(query_keys), 1)
            ranked.append((score, row))
        ranked.sort(key=lambda item: item[0], reverse=True)
        results: list[Evidence] = []
        seen: set[str] = set()
        for score, row in ranked:
            normalized = re.sub(r"\s+", "", row["content"]).lower()
            fingerprint = hashlib.sha256(normalized.encode("utf-8")).hexdigest()
            if score <= 0 or fingerprint in seen:
                continue
            seen.add(fingerprint)
            results.append(Evidence(
                row["original_name"], row["section"] or "正文", row["content"], round(score, 4),
                str(row.get("material_type") or ""), str(row.get("material_label") or ""),
            ))
            if len(results) >= max(1, min(int(top_k), 20)):
                break
        return results


class CampusService:
    def __init__(self, db: LearningDatabase, storage_dir: Path | str | None = None,
                 provider_factory=None):
        self.db = db
        self.storage_dir = Path(storage_dir or DATA_DIR / "uploads")
        self.storage_dir.mkdir(parents=True, exist_ok=True)
        self.provider_factory = provider_factory or build_backend_provider

    def get_course(self, course_id: str) -> dict:
        course = self.db.fetch_one("SELECT * FROM courses WHERE course_id=?", (course_id,))
        if not course:
            raise NotFound("课程不存在")
        return course

    def _can_access(self, course: dict, user_id: str, role: str) -> bool:
        if course["course_type"] == "personal_course":
            return role == "student" and course["owner_id"] == user_id
        if role == "teacher":
            return course["owner_id"] == user_id
        if role == "student":
            if course["visibility"] == "public":
                return True
            return self.db.fetch_one("SELECT 1 ok FROM course_enrollments WHERE course_id=? AND student_id=?",
                                     (course["course_id"], user_id)) is not None
        return False

    def require_access(self, course_id: str, user_id: str, role: str) -> dict:
        course = self.get_course(course_id)
        if not self._can_access(course, user_id, role):
            raise PermissionDenied("无权访问该课程")
        return course

    def list_courses(self, user_id: str, role: str) -> list[dict]:
        rows = self.db.fetch_all("SELECT * FROM courses ORDER BY updated_at DESC")
        visible = [row for row in rows if self._can_access(row, user_id, role)]
        for row in visible:
            published = self.db.fetch_one(
                """SELECT version_id FROM knowledge_versions WHERE course_id=? AND status='published'
                   ORDER BY version_number DESC LIMIT 1""", (row["course_id"],)
            )
            if published:
                materials = self.db.fetch_all(
                    """SELECT DISTINCT n.material_type FROM knowledge_version_nodes vn
                       JOIN knowledge_nodes n USING(node_id)
                       WHERE vn.version_id=? AND n.node_type='knowledge_point' AND n.status='approved'
                       ORDER BY n.material_type""", (published["version_id"],)
                )
            else:
                materials = self.db.fetch_all(
                    """SELECT DISTINCT COALESCE(m.material_type,'other') material_type
                       FROM course_documents d LEFT JOIN document_material_metadata m USING(document_id)
                       WHERE d.course_id=? AND d.status='ready' ORDER BY material_type""",
                    (row["course_id"],),
                )
            row["material_partitions"] = [{
                "material_type": str(item["material_type"]),
                "label": COURSE_MATERIAL_LABELS.get(str(item["material_type"]), "其他"),
            } for item in materials]
        return visible

    def create_course(self, name: str, course_type: str, user_id: str, role: str,
                      description: str = "", visibility: str | None = None) -> dict:
        name = name.strip()
        if not name:
            raise ValidationError("课程名称不能为空")
        expected = "personal_course" if role == "student" else "shared_course" if role == "teacher" else ""
        if course_type != expected:
            raise PermissionDenied("学生只能创建个人课程，教师只能创建共享课程")
        course_id = f"{'pc' if course_type == 'personal_course' else 'sc'}_{uuid.uuid4().hex[:12]}"
        self.db.execute("""INSERT INTO courses(course_id,course_name,course_type,owner_id,created_by_role,visibility,description)
                         VALUES(?,?,?,?,?,?,?)""",
                        (course_id, name, course_type, user_id, role,
                         visibility or ("private" if course_type == "personal_course" else "enrolled"), description.strip()))
        return self.get_course(course_id)

    def update_course(self, course_id: str, user_id: str, role: str, **changes: Any) -> dict:
        course = self.require_access(course_id, user_id, role)
        if course["owner_id"] != user_id:
            raise PermissionDenied("只有课程所有者可以修改课程")
        name = str(changes.get("course_name", course["course_name"])).strip()
        description = str(changes.get("description", course["description"])).strip()
        visibility = str(changes.get("visibility", course["visibility"]))
        if course["course_type"] == "personal_course":
            visibility = "private"
        if visibility not in {"private", "enrolled", "public"}:
            raise ValidationError("课程可见性不合法")
        self.db.execute("UPDATE courses SET course_name=?,description=?,visibility=?,updated_at=CURRENT_TIMESTAMP WHERE course_id=?",
                        (name, description, visibility, course_id))
        return self.get_course(course_id)

    def delete_personal_course(self, course_id: str, user_id: str) -> None:
        course = self.require_access(course_id, user_id, "student")
        if course["course_type"] != "personal_course" or course["owner_id"] != user_id:
            raise PermissionDenied("只能删除自己创建的个人课程")
        upload_root = self.storage_dir.resolve()
        course_dir = (upload_root / course_id).resolve()
        if upload_root not in course_dir.parents:
            raise ValidationError("课程存储路径不安全")
        self.db.execute("DELETE FROM courses WHERE course_id=?", (course_id,))
        if course_dir.exists():
            shutil.rmtree(course_dir)

    def upsert_virtual_course(self, course_id: str, name: str, teacher_id: str,
                              description: str = "", visibility: str = "public") -> dict:
        """Stable insertion/update hook for demo or externally provisioned courses."""
        if not re.fullmatch(r"[A-Za-z0-9_-]{3,64}", course_id):
            raise ValidationError("虚拟课程 ID 只能包含字母、数字、下划线和连字符")
        existing = self.db.fetch_one("SELECT owner_id FROM courses WHERE course_id=?", (course_id,))
        if existing and existing["owner_id"] != teacher_id:
            raise PermissionDenied("不能修改其他教师拥有的虚拟课程")
        self.db.execute("""INSERT INTO courses(course_id,course_name,course_type,owner_id,created_by_role,visibility,description,is_virtual)
                         VALUES(?,?,'shared_course',?,'teacher',?,?,1)
                         ON CONFLICT(course_id) DO UPDATE SET course_name=excluded.course_name,
                         owner_id=excluded.owner_id,visibility=excluded.visibility,description=excluded.description,
                         is_virtual=1,updated_at=CURRENT_TIMESTAMP""",
                        (course_id, name.strip(), teacher_id, visibility, description.strip()))
        return self.get_course(course_id)

    def enroll_student(self, course_id: str, teacher_id: str, student_id: str) -> None:
        course = self.require_access(course_id, teacher_id, "teacher")
        if course["course_type"] != "shared_course" or course["owner_id"] != teacher_id:
            raise PermissionDenied("只能为自己创建的共享课程授权")
        if not student_id.strip():
            raise ValidationError("学生 ID 不能为空")
        self.db.execute("INSERT OR IGNORE INTO course_enrollments(course_id,student_id) VALUES(?,?)",
                        (course_id, student_id.strip()))

    def list_enrollments(self, course_id: str, teacher_id: str) -> list[dict]:
        course = self.require_access(course_id, teacher_id, "teacher")
        if course["owner_id"] != teacher_id:
            raise PermissionDenied("无权查看授权名单")
        return self.db.fetch_all("SELECT student_id,created_at FROM course_enrollments WHERE course_id=?", (course_id,))

    def upload_document(self, course_id: str, user_id: str, role: str, name: str,
                        mime_type: str, data: bytes) -> dict:
        course = self.require_access(course_id, user_id, role)
        if role == "student" and (course["course_type"] != "personal_course" or course["owner_id"] != user_id):
            raise PermissionDenied("学生只能向自己的个人课程上传资料")
        if role == "teacher" and (course["course_type"] != "shared_course" or course["owner_id"] != user_id):
            raise PermissionDenied("教师只能向自己的共享课程上传资料")
        safe_name = _safe_name(name)
        suffix = Path(safe_name).suffix.lower()
        if suffix not in ALLOWED_FILES or mime_type not in ALLOWED_FILES[suffix]:
            raise ValidationError("扩展名与 MIME 类型不匹配或不受支持")
        try:
            validate_document_bytes(safe_name, data, max_bytes=MAX_UPLOAD_BYTES)
        except UnsafeUpload as exc:
            raise ValidationError(str(exc)) from exc
        digest = hashlib.sha256(data).hexdigest()
        if self.db.fetch_one("SELECT 1 ok FROM course_documents WHERE course_id=? AND sha256=?", (course_id, digest)):
            raise ValidationError("该课程中已存在内容相同的文件")
        document_id = f"doc_{uuid.uuid4().hex}"
        destination = (self.storage_dir / course_id / f"{document_id}_{safe_name}").resolve()
        root = (self.storage_dir / course_id).resolve()
        if root not in destination.parents:
            raise ValidationError("文件路径不安全")
        try:
            chunks = parse_document(data, suffix)
        except ValidationError:
            raise
        except Exception as exc:
            raise ValidationError("文件解析失败，请确认文件未损坏且格式正确") from exc
        if not chunks:
            raise ValidationError("文件未解析出有效文字")
        destination.parent.mkdir(parents=True, exist_ok=True)
        destination.write_bytes(data)
        with self.db.connect() as conn:
            conn.execute("""INSERT INTO course_documents(document_id,course_id,uploader_id,original_name,stored_path,mime_type,size_bytes,sha256,status)
                          VALUES(?,?,?,?,?,?,?,?,?)""",
                         (document_id, course_id, user_id, safe_name, str(destination), mime_type, len(data), digest, "ready"))
            conn.executemany("INSERT INTO document_chunks(document_id,course_id,section,page_number,content) VALUES(?,?,?,?,?)",
                             [(document_id, course_id, x["section"], x["page_number"], x["content"]) for x in chunks])
        return {"document_id": document_id, "file_name": safe_name, "status": "ready", "chunk_count": len(chunks)}

    def list_documents(self, course_id: str, user_id: str, role: str) -> list[dict]:
        self.require_access(course_id, user_id, role)
        return self.db.fetch_all("""SELECT d.document_id,d.original_name,d.mime_type,d.size_bytes,d.status,d.error_message,d.created_at,
                                   COUNT(c.chunk_id) chunk_count,
                                   (SELECT dc.content FROM document_chunks dc WHERE dc.document_id=d.document_id ORDER BY dc.chunk_id LIMIT 1) text_preview
                                   FROM course_documents d LEFT JOIN document_chunks c ON c.document_id=d.document_id
                                   WHERE d.course_id=? GROUP BY d.document_id ORDER BY d.created_at DESC""", (course_id,))

    def delete_document(self, document_id: str, user_id: str, role: str) -> None:
        doc = self.db.fetch_one("SELECT d.*,c.course_type,c.owner_id FROM course_documents d JOIN courses c USING(course_id) WHERE document_id=?",
                                (document_id,))
        if not doc:
            raise NotFound("资料不存在")
        if doc["owner_id"] != user_id or (role == "student" and doc["course_type"] != "personal_course") or (role == "teacher" and doc["course_type"] != "shared_course"):
            raise PermissionDenied("无权删除该资料")
        self.db.execute("DELETE FROM course_documents WHERE document_id=?", (document_id,))
        path = Path(doc["stored_path"])
        if path.exists() and self.storage_dir.resolve() in path.resolve().parents:
            path.unlink()

    def knowledge_status(self, course_id: str, user_id: str, role: str) -> dict:
        self.require_access(course_id, user_id, role)
        row = self.db.fetch_one("""SELECT COUNT(DISTINCT d.document_id) document_count,COUNT(c.chunk_id) chunk_count,
                                  COALESCE(SUM(LENGTH(c.content)),0) character_count
                                  FROM course_documents d LEFT JOIN document_chunks c ON c.document_id=d.document_id WHERE d.course_id=?""", (course_id,))
        return {**(row or {}), "status": "ready" if row and row["chunk_count"] else "empty"}

    def _retriever(self, course_id: str, material_type: str | None = None) -> ChunkRetriever:
        if material_type is not None and material_type not in COURSE_MATERIAL_LABELS:
            raise ValidationError("资料用途类型无效")
        published = self.db.fetch_one(
            "SELECT version_id FROM knowledge_versions WHERE course_id=? AND status='published' ORDER BY version_number DESC LIMIT 1",
            (course_id,),
        )
        if published:
            material_condition = ""
            params: tuple[Any, ...] = (published["version_id"],)
            if material_type:
                material_condition = " AND n.material_type=?"
                params += (material_type,)
            semantic_rows = self.db.fetch_all(
                """SELECT n.markdown content,
                          CASE n.material_type
                              WHEN 'slides' THEN '课件' WHEN 'textbook' THEN '教材'
                              WHEN 'syllabus' THEN '教学大纲' WHEN 'lesson_plan' THEN '教案'
                              WHEN 'experiment' THEN '实验资料' WHEN 'question_bank' THEN '题库'
                              WHEN 'knowledge_graph' THEN '知识图谱'
                              WHEN 'teaching_schedule' THEN '教学进度' ELSE '其他' END
                          || ' · ' || COALESCE(c.title,'')
                          || CASE WHEN s.title IS NOT NULL THEN ' / ' || s.title ELSE '' END section,
                          (SELECT MIN(ns.page_number) FROM knowledge_node_sources ns WHERE ns.node_id=n.node_id) page_number,
                          COALESCE((SELECT GROUP_CONCAT(DISTINCT d.original_name) FROM knowledge_node_sources ns
                                    JOIN course_documents d USING(document_id) WHERE ns.node_id=n.node_id),'课程知识库') original_name,
                          n.material_type,
                          CASE n.material_type
                              WHEN 'slides' THEN '课件' WHEN 'textbook' THEN '教材'
                              WHEN 'syllabus' THEN '教学大纲' WHEN 'lesson_plan' THEN '教案'
                              WHEN 'experiment' THEN '实验资料' WHEN 'question_bank' THEN '题库'
                              WHEN 'knowledge_graph' THEN '知识图谱'
                              WHEN 'teaching_schedule' THEN '教学进度' ELSE '其他' END material_label
                   FROM knowledge_version_nodes vn JOIN knowledge_nodes n USING(node_id)
                   LEFT JOIN knowledge_nodes s ON s.node_id=n.parent_id
                   LEFT JOIN knowledge_nodes c ON c.node_id=s.parent_id
                   WHERE vn.version_id=? AND n.status='approved' AND TRIM(n.markdown)<>''
                """ + material_condition + " ORDER BY n.sort_order", params,
            )
            if semantic_rows:
                return ChunkRetriever(semantic_rows)
            block_condition = ""
            block_params: tuple[Any, ...] = (published["version_id"],)
            if material_type:
                block_condition = " AND COALESCE(m.material_type,'other')=?"
                block_params += (material_type,)
            rows = self.db.fetch_all(
                """SELECT b.plain_text content,'第 ' || COALESCE(b.page_number,1) || ' 页' section,b.page_number,d.original_name
                          ,COALESCE(m.material_type,'other') material_type
                   FROM knowledge_version_blocks vb JOIN document_blocks b USING(block_id)
                   JOIN course_documents d USING(document_id)
                   LEFT JOIN document_material_metadata m USING(document_id)
                   WHERE vb.version_id=? AND b.visibility_level='PUBLIC'
                     AND b.verification_status IN ('auto_verified','teacher_verified')"""
                + block_condition, block_params,
            )
            for row in rows:
                row["material_label"] = COURSE_MATERIAL_LABELS.get(row["material_type"], "其他")
                row["section"] = f"{row['material_label']} · {row['section']}"
            return ChunkRetriever(rows)
        chunk_condition = ""
        chunk_params: tuple[Any, ...] = (course_id,)
        if material_type:
            chunk_condition = " AND COALESCE(m.material_type,'other')=?"
            chunk_params += (material_type,)
        rows = self.db.fetch_all("""SELECT c.content,c.section,c.page_number,d.original_name,
                                         COALESCE(m.material_type,'other') material_type
                                  FROM document_chunks c JOIN course_documents d USING(document_id)
                                  LEFT JOIN document_material_metadata m USING(document_id)
                                  WHERE c.course_id=? AND d.status='ready'""" + chunk_condition, chunk_params)
        for row in rows:
            row["material_label"] = COURSE_MATERIAL_LABELS.get(row["material_type"], "其他")
            row["section"] = f"{row['material_label']} · {row['section']}"
        return ChunkRetriever(rows)

    def ask(self, course_id: str, user_id: str, role: str, question: str,
            provider: LLMProvider | None = None, *, intent: str | None = None,
            student_message: str = "", phase: str = "initial",
            history: list[dict[str, Any]] | None = None,
            evidence_refs: list[dict[str, Any]] | None = None,
            retrieval_scope: str = "all", material_type: str | None = None,
            session_id: str | None = None) -> dict:
        self.require_access(course_id, user_id, role)
        clean_question = question.strip()
        if not clean_question:
            raise ValidationError("问题不能为空")
        if retrieval_scope not in {"all", "material"}:
            raise ValidationError("检索范围必须是 all 或 material")
        if retrieval_scope == "material" and material_type not in COURSE_MATERIAL_LABELS:
            raise ValidationError("指定材料检索必须提供有效 material_type")
        selected_material = material_type if retrieval_scope == "material" else None
        active_provider = provider or self.provider_factory()
        if intent is not None:
            if role != "student":
                raise PermissionDenied("引导式问答仅供学生使用")
            stored = None
            if intent == "start":
                phase = "initial"
                history = []
                evidence_refs = []
            else:
                if not session_id:
                    raise ValidationError("引导会话已失效，请重新开始本题")
                stored = self.db.fetch_one(
                    "SELECT * FROM guided_qa_sessions WHERE session_id=?", (session_id,),
                )
                if not stored or stored["user_id"] != user_id or stored["course_id"] != course_id:
                    raise PermissionDenied("无权访问该引导会话")
                if stored["status"] != "active":
                    raise ValidationError("本题会话已结束，请重新开始")
                if stored["question"] != clean_question:
                    raise ValidationError("问题与当前引导会话不一致")
                phase = stored["phase"]
                history = _loads(stored["history_json"])
                evidence_refs = _loads(stored["evidence_json"])
                retrieval_scope = stored["retrieval_scope"] or "all"
                material_type = stored["material_type"]
                selected_material = material_type if retrieval_scope == "material" else None
                if intent == "reveal" and int(stored["turn_count"]) < 2:
                    message = "我先不直接给答案。请至少写出两步自己的判断；如果卡住，可以先要一个提示。"
                    return {
                        "session_id": session_id, "question_id": None,
                        "reply": message, "answer": message,
                        "phase": "guiding", "expects_response": True, "can_reveal": False,
                        "completed": False, "sources": evidence_refs,
                        "knowledge_points": [], "refused": False, "persisted": False,
                    }
            retriever = self._retriever(course_id, selected_material)
            result = guide_question(
                clean_question, retriever, active_provider,
                intent=intent, phase=phase, student_message=student_message,
                history=history, evidence_refs=evidence_refs,
                min_score=MIN_EVIDENCE_SCORE, top_k=TOP_K,
            )
            history_rows = list(history or [])
            next_turn_count = 0
            if intent == "start":
                session_id = f"gqa_{uuid.uuid4().hex}"
                history_rows = [
                    {"role": "user", "content": clean_question},
                    {"role": "assistant", "content": result.reply},
                ]
                self.db.execute(
                    """INSERT INTO guided_qa_sessions(
                           session_id,course_id,user_id,question,phase,turn_count,
                           evidence_json,history_json,retrieval_scope,material_type,status
                       ) VALUES(?,?,?,?,?,0,?,?,?,?,?)""",
                    (session_id, course_id, user_id, clean_question, result.phase,
                     _json([item.to_dict() for item in result.evidence]), _json(history_rows),
                     retrieval_scope, selected_material,
                     "closed" if result.completed else "active"),
                )
            elif stored:
                if student_message.strip():
                    history_rows.append({"role": "user", "content": student_message.strip()[:2000]})
                history_rows.append({"role": "assistant", "content": result.reply})
                next_turn_count = int(stored["turn_count"]) + (1 if intent == "respond" else 0)
                next_status = "revealed" if intent == "reveal" and not result.refused else (
                    "closed" if result.completed else "active"
                )
                self.db.execute(
                    """UPDATE guided_qa_sessions SET phase=?,turn_count=?,history_json=?,status=?,
                           updated_at=CURRENT_TIMESTAMP WHERE session_id=?""",
                    (result.phase, next_turn_count, _json(history_rows[-12:]), next_status, session_id),
                )
            should_persist = result.refused or (intent == "reveal" and not result.refused)
            question_id = None
            if should_persist:
                question_id = self.db.execute(
                    """INSERT INTO course_questions(course_id,user_id,question,answer,sources_json,knowledge_points_json,refused)
                       VALUES(?,?,?,?,?,?,?)""",
                    (
                        course_id, user_id, clean_question, result.reply,
                        _json([item.to_dict() for item in result.evidence]),
                        _json(result.knowledge_points), int(result.refused),
                    ),
                )
            return {
                "session_id": session_id,
                "question_id": question_id,
                "reply": result.reply,
                "answer": result.reply,
                "phase": result.phase,
                "expects_response": result.expects_response,
                "can_reveal": bool(result.can_reveal and next_turn_count >= 2),
                "completed": result.completed,
                "sources": [item.to_dict() for item in result.evidence],
                "knowledge_points": result.knowledge_points,
                "refused": result.refused,
                "persisted": should_persist,
            }
        retriever = self._retriever(course_id, selected_material)
        result = answer_question(clean_question, retriever, active_provider,
                                 MIN_EVIDENCE_SCORE, TOP_K)
        qid = self.db.execute("""INSERT INTO course_questions(course_id,user_id,question,answer,sources_json,knowledge_points_json,refused)
                               VALUES(?,?,?,?,?,?,?)""",
                              (course_id,user_id,clean_question,result.answer,_json([e.to_dict() for e in result.evidence]),
                               _json(result.knowledge_points),int(result.refused)))
        return {"question_id": qid, "answer": result.answer, "sources": [e.to_dict() for e in result.evidence],
                "knowledge_points": result.knowledge_points, "refused": result.refused}

    def generate_quiz(self, course_id: str, user_id: str, role: str, question_id: int | None = None) -> dict:
        self.require_access(course_id, user_id, role)
        params: tuple = (course_id, user_id)
        sql = "SELECT * FROM course_questions WHERE course_id=? AND user_id=? AND refused=0"
        if question_id is not None:
            sql += " AND question_id=?"
            params += (question_id,)
        row = self.db.fetch_one(sql + " ORDER BY question_id DESC LIMIT 1", params)
        if not row:
            raise ValidationError("请先完成一次有资料证据支持的课程问答")
        weak = self.profile(course_id, user_id, role)["weak_points"]
        items = generate_exercises(row["answer"], _loads(row["knowledge_points_json"]), weak)
        return {"question_id": row["question_id"], "items": [item.to_dict() for item in items]}

    def submit_quiz(self, course_id: str, user_id: str, role: str, question_id: int | None,
                    items: list[dict], responses: list[str | None]) -> dict:
        self.require_access(course_id, user_id, role)
        exercise_items = [ExerciseItem(**item) for item in items]
        grade = grade_exercises(exercise_items, responses)
        points = sorted({point for item in exercise_items for point in item.knowledge_points})
        attempt_id = self.db.execute("""INSERT INTO course_attempts(course_id,user_id,question_id,score,total,wrong_items_json,records_json,knowledge_points_json)
                                    VALUES(?,?,?,?,?,?,?,?)""",
                                   (course_id,user_id,question_id,grade.score,grade.total,_json(grade.wrong_items),
                                    _json(grade.records),_json(points)))
        self.db.update_course_points(course_id, user_id, grade.records)
        return {"attempt_id": attempt_id, "score": grade.score, "correct_count": grade.correct_count,
                "total": grade.total, "records": grade.records, "wrong_items": grade.wrong_items,
                "topic_stats": grade.topic_stats}

    def profile(self, course_id: str, user_id: str, role: str) -> dict:
        self.require_access(course_id, user_id, role)
        if role != "student":
            raise PermissionDenied("学习画像仅供学生本人查看")
        questions = self.db.fetch_all("SELECT * FROM course_questions WHERE course_id=? AND user_id=? ORDER BY question_id DESC LIMIT 20", (course_id,user_id))
        attempts = self.db.fetch_all("SELECT * FROM course_attempts WHERE course_id=? AND user_id=? ORDER BY attempt_id DESC LIMIT 20", (course_id,user_id))
        weak = self.db.fetch_all("SELECT * FROM course_weak_points WHERE course_id=? AND user_id=? ORDER BY weakness_score DESC,answered DESC", (course_id,user_id))
        for row in questions:
            row["sources"] = _loads(row.pop("sources_json")); row["knowledge_points"] = _loads(row.pop("knowledge_points_json"))
        for row in attempts:
            row["wrong_items"] = _loads(row.pop("wrong_items_json")); row["records"] = _loads(row.pop("records_json")); row["knowledge_points"] = _loads(row.pop("knowledge_points_json"))
        return {"questions": questions, "attempts": attempts, "weak_points": weak,
                "wrong_questions": [x for a in attempts for x in a["wrong_items"]]}

    def class_analysis(self, course_id: str, teacher_id: str) -> dict:
        course = self.require_access(course_id, teacher_id, "teacher")
        if course["course_type"] != "shared_course" or course["owner_id"] != teacher_id:
            raise PermissionDenied("教师只能分析自己的共享课程")
        questions = self.db.fetch_all("SELECT question,refused,knowledge_points_json FROM course_questions WHERE course_id=?", (course_id,))
        attempts = self.db.fetch_all("SELECT score,total,records_json FROM course_attempts WHERE course_id=?", (course_id,))
        frequent = [{"question": q, "count": n} for q,n in Counter(x["question"] for x in questions).most_common(10)]
        uncovered = [{"question": q, "count": n} for q,n in Counter(x["question"] for x in questions if x["refused"]).most_common(10)]
        point_stats: dict[str, dict[str,int]] = defaultdict(lambda: {"answered":0,"correct":0})
        for attempt in attempts:
            for record in _loads(attempt["records_json"]):
                for point in record.get("knowledge_points", []):
                    point_stats[point]["answered"] += 1
                    point_stats[point]["correct"] += int(bool(record.get("correct")))
        weak = [{"knowledge_point": p, **s, "accuracy": round(100*s["correct"]/s["answered"],1) if s["answered"] else 0}
                for p,s in point_stats.items()]
        weak.sort(key=lambda x: (x["accuracy"], -x["answered"]))
        avg = round(sum(x["score"] for x in attempts)/len(attempts),1) if attempts else 0
        return {"question_count": len(questions), "quiz_count": len(attempts), "average_score": avg,
                "frequent_questions": frequent, "uncovered_questions": uncovered, "weak_points": weak,
                "privacy": "仅展示共享课程的匿名聚合结果"}

    def teaching_report(self, course_id: str, teacher_id: str) -> dict:
        course = self.get_course(course_id)
        analysis = self.class_analysis(course_id, teacher_id)
        phenomena, evidence, suggestions = [], [], []
        if analysis["weak_points"]:
            point = analysis["weak_points"][0]
            phenomena.append(f"学生在“{point['knowledge_point']}”上的掌握较弱")
            evidence.append(f"该知识点共作答 {point['answered']} 次，正确率 {point['accuracy']}%")
            suggestions.append(f"围绕“{point['knowledge_point']}”增加例题讲解与分层练习")
        if analysis["uncovered_questions"]:
            phenomena.append("部分学生问题未被现有资料充分覆盖")
            evidence.append(f"共有 {sum(x['count'] for x in analysis['uncovered_questions'])} 次问题因资料不足未能回答")
            suggestions.append("补充对应章节资料，并在入库后重新检查检索覆盖率")
        if not phenomena:
            phenomena.append("当前共享课程尚缺少足够学习行为数据")
            evidence.append(f"已记录 {analysis['question_count']} 次提问和 {analysis['quiz_count']} 次练习")
            suggestions.append("先组织一次课程问答与形成性练习，再依据真实数据调整教学")
        return {"course_id": course_id, "course_name": course["course_name"], "phenomena": phenomena,
                "evidence": evidence, "suggestions": suggestions, "analysis": analysis}

    def export_student_csv(self, course_id: str, user_id: str) -> bytes:
        profile = self.profile(course_id, user_id, "student")
        output = io.StringIO(); writer = csv.writer(output)
        writer.writerow(["类型","时间","内容","成绩/状态"])
        for row in profile["questions"]:
            writer.writerow(["提问",row["created_at"],row["question"],"资料不足" if row["refused"] else "已回答"])
        for row in profile["attempts"]:
            writer.writerow(["练习",row["created_at"],f"错题 {len(row['wrong_items'])} 道",row["score"]])
        return ("\ufeff"+output.getvalue()).encode("utf-8")

    def export_class_csv(self, course_id: str, teacher_id: str) -> bytes:
        report = self.teaching_report(course_id, teacher_id)
        output = io.StringIO(); writer = csv.writer(output)
        writer.writerow(["现象","证据","建议"])
        for i in range(max(len(report["phenomena"]),len(report["evidence"]),len(report["suggestions"]))):
            writer.writerow([report["phenomena"][i] if i < len(report["phenomena"]) else "",
                             report["evidence"][i] if i < len(report["evidence"]) else "",
                             report["suggestions"][i] if i < len(report["suggestions"]) else ""])
        return ("\ufeff"+output.getvalue()).encode("utf-8")

    def export_class_excel(self, course_id: str, teacher_id: str) -> bytes:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
        report = self.teaching_report(course_id, teacher_id)
        workbook = Workbook(); sheet = workbook.active; sheet.title = "教学改进"
        sheet.append(["现象", "证据", "建议"])
        for cell in sheet[1]:
            cell.font = Font(bold=True, color="FFFFFF"); cell.fill = PatternFill("solid", fgColor="D98686")
        for i in range(max(len(report["phenomena"]), len(report["evidence"]), len(report["suggestions"]))):
            sheet.append([report["phenomena"][i] if i < len(report["phenomena"]) else "",
                          report["evidence"][i] if i < len(report["evidence"]) else "",
                          report["suggestions"][i] if i < len(report["suggestions"]) else ""])
        sheet.column_dimensions["A"].width = 34; sheet.column_dimensions["B"].width = 42; sheet.column_dimensions["C"].width = 42
        analysis = workbook.create_sheet("匿名聚合")
        analysis.append(["指标", "数值"])
        analysis.append(["提问次数", report["analysis"]["question_count"]])
        analysis.append(["练习次数", report["analysis"]["quiz_count"]])
        analysis.append(["平均分", report["analysis"]["average_score"]])
        output = io.BytesIO(); workbook.save(output); return output.getvalue()

    def export_class_word(self, course_id: str, teacher_id: str) -> bytes:
        from docx import Document
        from docx.shared import Pt
        report = self.teaching_report(course_id, teacher_id)
        document = Document(); document.add_heading(f"{report['course_name']} 学情报告", 0)
        document.add_paragraph("本报告仅使用共享课程中的匿名聚合数据。")
        for index, phenomenon in enumerate(report["phenomena"]):
            document.add_heading(f"改进建议 {index + 1}", level=1)
            document.add_paragraph(f"现象：{phenomenon}")
            document.add_paragraph(f"证据：{report['evidence'][min(index, len(report['evidence'])-1)]}")
            document.add_paragraph(f"建议：{report['suggestions'][min(index, len(report['suggestions'])-1)]}")
        style = document.styles["Normal"]; style.font.name = "Microsoft YaHei"; style.font.size = Pt(10.5)
        output = io.BytesIO(); document.save(output); return output.getvalue()

    def seed_demo(self, materials_dir: Path) -> None:
        self.upsert_virtual_course("virtual_ai_101", "人工智能基础（虚拟课程）", "demo_teacher_001",
                                   "用于演示教师共享课程、问答、练习和匿名学情分析。", "public")
        for path in sorted(materials_dir.glob("*.md")):
            data = path.read_bytes()
            document = self.db.fetch_one(
                "SELECT * FROM course_documents WHERE course_id='virtual_ai_101' AND original_name=?",
                (path.name,),
            )
            if not document:
                self.upload_document("virtual_ai_101", "demo_teacher_001", "teacher", path.name, "text/markdown", data)
                continue
            stored = Path(str(document.get("stored_path") or ""))
            if not stored.is_file():
                destination = (self.storage_dir / "virtual_ai_101" / f"{document['document_id']}_{_safe_name(path.name)}").resolve()
                destination.parent.mkdir(parents=True, exist_ok=True)
                destination.write_bytes(data)
                self.db.execute(
                    "UPDATE course_documents SET stored_path=?,size_bytes=?,sha256=?,status='ready',error_message='' WHERE document_id=?",
                    (str(destination), len(data), hashlib.sha256(data).hexdigest(), document["document_id"]),
                )
            chunks = self.db.fetch_one(
                "SELECT COUNT(*) n FROM document_chunks WHERE document_id=?", (document["document_id"],)
            )
            if not chunks or not chunks["n"]:
                parsed = parse_document(data, ".md")
                with self.db.connect() as conn:
                    conn.executemany(
                        "INSERT INTO document_chunks(document_id,course_id,section,page_number,content) VALUES(?,?,?,?,?)",
                        [(document["document_id"], "virtual_ai_101", item["section"], item["page_number"], item["content"]) for item in parsed],
                    )
        self._publish_demo_knowledge()

    def _publish_demo_knowledge(self) -> None:
        """Create a small, source-bound published snapshot for reliable demos."""
        existing = self.db.fetch_one(
            """SELECT 1 ok FROM knowledge_versions v JOIN knowledge_version_nodes n USING(version_id)
               WHERE v.course_id='virtual_ai_101' AND v.status='published' LIMIT 1"""
        )
        if existing:
            return
        chunks = self.db.fetch_all(
            """SELECT c.chunk_id,c.document_id,c.section,c.page_number,c.content,d.original_name
               FROM document_chunks c JOIN course_documents d USING(document_id)
               WHERE c.course_id='virtual_ai_101' AND d.status='ready'
               ORDER BY d.original_name,c.chunk_id LIMIT 24"""
        )
        if not chunks:
            return
        last = self.db.fetch_one(
            "SELECT COALESCE(MAX(version_number),0) n FROM knowledge_versions WHERE course_id='virtual_ai_101'"
        )
        version_number = int((last or {}).get("n") or 0) + 1
        version_id = f"demo_version_{version_number}"
        with self.db.connect() as conn:
            conn.execute(
                "UPDATE knowledge_versions SET status='superseded' WHERE course_id='virtual_ai_101' AND status='published'"
            )
            conn.execute(
                """INSERT INTO knowledge_versions(version_id,course_id,version_number,status,created_by,published_at,markdown_snapshot)
                   VALUES(?,'virtual_ai_101',?,'published','demo_teacher_001',CURRENT_TIMESTAMP,?)""",
                (version_id, version_number, "\n\n".join(item["content"] for item in chunks)),
            )
            for order, item in enumerate(chunks, 1):
                fingerprint = hashlib.sha256(f"{item['document_id']}:{item['chunk_id']}".encode()).hexdigest()[:20]
                block_id = f"demo_block_{fingerprint}"
                node_id = f"demo_node_{fingerprint}"
                page_number = int(item.get("page_number") or 1)
                title = str(item.get("section") or item["original_name"] or f"知识点 {order}")[:160]
                conn.execute(
                    """INSERT OR IGNORE INTO document_blocks(
                           block_id,document_id,block_order,block_type,markdown,plain_text,page_number,
                           verification_status,visibility_level,source_method,parser_name,parser_version,
                           reviewed_by,reviewed_at)
                       VALUES(?,?,?,'paragraph',?,?,?,'teacher_verified','PUBLIC','demo_seed',
                              'demo_seed','1','demo_teacher_001',CURRENT_TIMESTAMP)""",
                    (block_id, item["document_id"], order, item["content"], item["content"], page_number),
                )
                conn.execute(
                    """INSERT OR IGNORE INTO knowledge_nodes(
                           node_id,course_id,document_id,node_scope,node_type,title,summary,markdown,
                           keywords_json,source_pages_json,sort_order,status,reviewed_by,reviewed_at)
                       VALUES(?,'virtual_ai_101',?,'course','knowledge_point',?,?,?,?,?,?,'approved',
                              'demo_teacher_001',CURRENT_TIMESTAMP)""",
                    (node_id, item["document_id"], title, title, item["content"], "[]", json.dumps([page_number]), order),
                )
                conn.execute(
                    """INSERT OR IGNORE INTO knowledge_node_sources(node_id,block_id,document_id,page_number,bbox_json)
                       VALUES(?,?,?,?, '[]')""",
                    (node_id, block_id, item["document_id"], page_number),
                )
                conn.execute("INSERT OR IGNORE INTO knowledge_version_nodes(version_id,node_id) VALUES(?,?)", (version_id, node_id))
                conn.execute("INSERT OR IGNORE INTO knowledge_version_blocks(version_id,block_id) VALUES(?,?)", (version_id, block_id))
